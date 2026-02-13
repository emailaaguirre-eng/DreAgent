"""
=============================================================================
HUMMINGBIRD-LEA - PowerPoint Generation Service
Powered by CoDre-X | B & D Servicing LLC
=============================================================================
PowerPoint document generation using python-pptx.

Features:
- Professional slide layouts
- EIAG-branded presentations
- Tables, charts, and bullet points
- Template-based generation
=============================================================================
"""

import logging
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple
import time

from .base import (
    BaseDocumentGenerator,
    DocumentFormat,
    DocumentMetadata,
    DocumentSection,
    DocumentResult,
    ColorScheme,
    FontScheme,
    ProposalData,
    SiteSelectionData,
    IncentiveSummaryData,
)

logger = logging.getLogger(__name__)


# =============================================================================
# PowerPoint Generator
# =============================================================================

class PowerPointGenerator(BaseDocumentGenerator):
    """
    PowerPoint document generator using python-pptx.

    Usage:
        generator = PowerPointGenerator()

        # Generate from sections
        result = generator.generate(metadata, sections)

        # Generate EIAG proposal presentation
        result = generator.generate_proposal(proposal_data)
    """

    @property
    def format(self) -> DocumentFormat:
        return DocumentFormat.PPTX

    def __init__(
        self,
        output_dir: Optional[Path] = None,
        colors: Optional[ColorScheme] = None,
        fonts: Optional[FontScheme] = None,
    ):
        super().__init__(output_dir, colors, fonts)

        # Check for pptx availability
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            self._pptx_available = True
        except ImportError:
            self._pptx_available = False
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")

    def _check_pptx(self):
        """Check if python-pptx is available"""
        if not self._pptx_available:
            raise ImportError(
                "python-pptx is required for PowerPoint generation. "
                "Install with: pip install python-pptx"
            )

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def generate(
        self,
        metadata: DocumentMetadata,
        sections: List[DocumentSection],
    ) -> DocumentResult:
        """
        Generate a PowerPoint presentation from sections.

        Args:
            metadata: Document metadata
            sections: List of sections (each becomes a slide)

        Returns:
            DocumentResult with file path
        """
        self._check_pptx()
        start_time = time.time()

        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN

        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
            prs.slide_height = Inches(7.5)

            # Add title slide
            self._add_title_slide(prs, metadata)

            # Add content slides
            for section in sections:
                self._add_content_slide(prs, section)

            # Add closing slide
            self._add_closing_slide(prs, metadata)

            # Save file
            filename = self._generate_filename(metadata)
            output_path = self._get_output_path(filename)
            prs.save(str(output_path))

            processing_time = (time.time() - start_time) * 1000

            return DocumentResult(
                success=True,
                file_path=output_path,
                format=self.format,
                file_size_bytes=output_path.stat().st_size,
                page_count=len(prs.slides),
                processing_time_ms=processing_time,
                metadata=metadata,
            )

        except Exception as e:
            logger.error(f"PowerPoint generation error: {e}")
            return DocumentResult(
                success=False,
                error=str(e),
            )

    def _add_title_slide(self, prs, metadata: DocumentMetadata):
        """Add a title slide"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN

        # Use blank layout and build custom
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add background shape (optional colored bar)
        self._add_header_bar(slide)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = metadata.title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.name = self.fonts.heading
        title_para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.primary))
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle (project/client name)
        if metadata.client_name or metadata.project_name:
            subtitle_text = metadata.project_name or metadata.client_name
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle_text
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.name = self.fonts.body
            subtitle_para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.secondary))
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Add date and company
        date_text = metadata.created_at.strftime("%B %d, %Y")
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.paragraphs[0]
        footer_para.text = f"{metadata.company} | {date_text}"
        footer_para.font.size = Pt(14)
        footer_para.font.name = self.fonts.body
        footer_para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.text_light))
        footer_para.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs, section: DocumentSection):
        """Add a content slide for a section"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add header bar
        self._add_header_bar(slide, height=0.8)

        # Add section title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = section.title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.name = self.fonts.heading
        title_para.font.color.rgb = RgbColor(255, 255, 255)

        # Add content based on section type
        content_top = Inches(1.2)

        if section.table_data:
            self._add_table(slide, section.table_data, content_top)
        elif section.items:
            self._add_bullet_list(slide, section.items, content_top, section.content)
        else:
            self._add_text_content(slide, section.content, content_top)

        # Add subsections if present
        if section.subsections:
            for subsection in section.subsections:
                self._add_content_slide(prs, subsection)

    def _add_header_bar(self, slide, height: float = 1.0):
        """Add a colored header bar"""
        from pptx.util import Inches
        from pptx.dml.color import RgbColor

        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(13.333), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(*self._hex_to_rgb(self.colors.primary))
        shape.line.fill.background()

    def _add_text_content(self, slide, content: str, top: float):
        """Add text content to a slide"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor

        content_box = slide.shapes.add_textbox(
            Inches(0.75), top, Inches(11.833), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        para = content_frame.paragraphs[0]
        para.text = content
        para.font.size = Pt(18)
        para.font.name = self.fonts.body
        para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.text))

    def _add_bullet_list(
        self,
        slide,
        items: List[str],
        top: float,
        intro_text: Optional[str] = None,
    ):
        """Add a bullet list to a slide"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor

        content_box = slide.shapes.add_textbox(
            Inches(0.75), top, Inches(11.833), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Add intro text if provided
        if intro_text:
            para = content_frame.paragraphs[0]
            para.text = intro_text
            para.font.size = Pt(16)
            para.font.name = self.fonts.body
            para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.text))
            para.space_after = Pt(12)

        # Add bullet items
        for i, item in enumerate(items):
            if intro_text or i > 0:
                para = content_frame.add_paragraph()
            else:
                para = content_frame.paragraphs[0]

            para.text = f"• {item}"
            para.font.size = Pt(16)
            para.font.name = self.fonts.body
            para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.text))
            para.space_before = Pt(6)
            para.level = 0

    def _add_table(self, slide, table_data: List[List[str]], top: float):
        """Add a table to a slide"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor

        if not table_data:
            return

        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0

        if rows == 0 or cols == 0:
            return

        # Calculate dimensions
        table_width = Inches(11.833)
        table_height = Inches(min(rows * 0.5, 5.0))
        col_width = table_width / cols

        # Add table
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.75), top,
            table_width, table_height
        ).table

        # Style header row
        for col_idx, cell_text in enumerate(table_data[0]):
            cell = table.cell(0, col_idx)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(*self._hex_to_rgb(self.colors.primary))

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RgbColor(255, 255, 255)
            para.font.name = self.fonts.body

        # Fill data rows
        for row_idx, row_data in enumerate(table_data[1:], start=1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)

                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RgbColor(*self._hex_to_rgb(self.colors.background_alt))

                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.font.name = self.fonts.body
                para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.text))

    def _add_closing_slide(self, prs, metadata: DocumentMetadata):
        """Add a closing/thank you slide"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add header bar
        self._add_header_bar(slide)

        # Add thank you text
        thank_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        thank_frame = thank_box.text_frame
        thank_para = thank_frame.paragraphs[0]
        thank_para.text = "Thank You"
        thank_para.font.size = Pt(44)
        thank_para.font.bold = True
        thank_para.font.name = self.fonts.heading
        thank_para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.primary))
        thank_para.alignment = PP_ALIGN.CENTER

        # Add contact info
        contact_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5)
        )
        contact_frame = contact_box.text_frame
        contact_para = contact_frame.paragraphs[0]
        contact_para.text = f"{metadata.company}\n{metadata.author}"
        contact_para.font.size = Pt(18)
        contact_para.font.name = self.fonts.body
        contact_para.font.color.rgb = RgbColor(*self._hex_to_rgb(self.colors.text_light))
        contact_para.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # EIAG Template Methods
    # =========================================================================

    def generate_proposal(self, data: ProposalData) -> DocumentResult:
        """
        Generate an EIAG incentive proposal presentation.

        Args:
            data: Proposal data

        Returns:
            DocumentResult
        """
        metadata = DocumentMetadata(
            title=f"Incentive Proposal - {data.project_name}",
            client_name=data.client_name,
            project_name=data.project_name,
        )

        sections = [
            DocumentSection(
                title="Executive Summary",
                content=data.executive_summary,
            ),
            DocumentSection(
                title="Project Background",
                content=data.background,
            ),
            DocumentSection(
                title="Proposed Incentives",
                content="",
                table_data=[
                    ["Incentive", "Value", "Type", "Timeline"],
                    *[
                        [inc.get("name", ""), inc.get("value", ""), inc.get("type", ""), inc.get("timeline", "")]
                        for inc in data.proposed_incentives
                    ]
                ],
            ),
            DocumentSection(
                title="Financial Summary",
                content="",
                items=[
                    f"Total Incentive Value: {data.financial_summary.get('total_value', 'TBD')}",
                    f"Investment Required: {data.financial_summary.get('investment', 'TBD')}",
                    f"Net Benefit: {data.financial_summary.get('net_benefit', 'TBD')}",
                    f"ROI: {data.financial_summary.get('roi', 'TBD')}",
                ],
            ),
            DocumentSection(
                title="Implementation Timeline",
                content="",
                table_data=[
                    ["Phase", "Timeline", "Milestone"],
                    *[
                        [item.get("phase", ""), item.get("timeline", ""), item.get("milestone", "")]
                        for item in data.timeline
                    ]
                ],
            ),
            DocumentSection(
                title="Next Steps",
                content="",
                items=data.next_steps,
            ),
        ]

        return self.generate(metadata, sections)

    def generate_site_selection(self, data: SiteSelectionData) -> DocumentResult:
        """
        Generate a site selection report presentation.

        Args:
            data: Site selection data

        Returns:
            DocumentResult
        """
        metadata = DocumentMetadata(
            title=f"Site Selection Report - {data.project_name}",
            client_name=data.client_name,
            project_name=data.project_name,
        )

        sections = [
            DocumentSection(
                title="Executive Summary",
                content=data.executive_summary,
            ),
            DocumentSection(
                title="Evaluation Criteria",
                content="",
                table_data=[
                    ["Criterion", "Weight", "Description"],
                    *[
                        [c.get("name", ""), c.get("weight", ""), c.get("description", "")]
                        for c in data.evaluation_criteria
                    ]
                ],
            ),
            DocumentSection(
                title="Sites Evaluated",
                content="",
                table_data=[
                    ["Site", "Location", "Size", "Key Features"],
                    *[
                        [s.get("name", ""), s.get("location", ""), s.get("size", ""), s.get("features", "")]
                        for s in data.sites_evaluated
                    ]
                ],
            ),
            DocumentSection(
                title="Comparison Matrix",
                content="",
                table_data=data.comparison_matrix,
            ),
            DocumentSection(
                title="Recommendation",
                content=data.recommendation,
            ),
        ]

        return self.generate(metadata, sections)

    def generate_client_presentation(
        self,
        title: str,
        client_name: str,
        agenda: List[str],
        key_points: List[Dict[str, Any]],
        recommendations: List[str],
        next_steps: List[str],
    ) -> DocumentResult:
        """
        Generate a general client presentation.

        Args:
            title: Presentation title
            client_name: Client name
            agenda: Agenda items
            key_points: List of key points with title and content
            recommendations: List of recommendations
            next_steps: List of next steps

        Returns:
            DocumentResult
        """
        metadata = DocumentMetadata(
            title=title,
            client_name=client_name,
        )

        sections = [
            DocumentSection(
                title="Agenda",
                content="",
                items=agenda,
            ),
        ]

        # Add key point slides
        for point in key_points:
            sections.append(DocumentSection(
                title=point.get("title", "Key Point"),
                content=point.get("content", ""),
                items=point.get("items"),
                table_data=point.get("table_data"),
            ))

        sections.extend([
            DocumentSection(
                title="Recommendations",
                content="",
                items=recommendations,
            ),
            DocumentSection(
                title="Next Steps",
                content="",
                items=next_steps,
            ),
        ])

        return self.generate(metadata, sections)


# =============================================================================
# Factory Function
# =============================================================================

_pptx_generator: Optional[PowerPointGenerator] = None


def get_powerpoint_generator() -> PowerPointGenerator:
    """Get or create the PowerPoint generator singleton"""
    global _pptx_generator
    if _pptx_generator is None:
        _pptx_generator = PowerPointGenerator()
    return _pptx_generator
