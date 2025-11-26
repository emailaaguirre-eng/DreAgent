"""
Lea Agentic Task System
Allows Lea to autonomously perform pre-configured tasks
"""

import os
import sys
import json
import logging
import subprocess
import shutil
import traceback
import html
import re
from pathlib import Path
from typing import Dict, List, Any, Callable, Optional, Tuple
from datetime import datetime, timedelta
from abc import ABC, abstractmethod

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Reports directory - all Lea-generated reports saved here
REPORTS_DIR = Path(r"F:\Dre_Programs\LeaAssistant\Lea_Created_Reports")
REPORTS_DIR.mkdir(parents=True, exist_ok=True)


def cleanup_old_reports(reports_dir: Path, days_old: int = 10):
    """Delete reports older than specified days"""
    try:
        cutoff_date = datetime.now() - timedelta(days=days_old)
        deleted_count = 0
        
        if not reports_dir.exists():
            return deleted_count
        
        for report_file in reports_dir.iterdir():
            if report_file.is_file():
                try:
                    # Get file modification time
                    file_mtime = datetime.fromtimestamp(report_file.stat().st_mtime)
                    if file_mtime < cutoff_date:
                        report_file.unlink()
                        deleted_count += 1
                        logger.info(f"Deleted old report: {report_file.name} (age: {(datetime.now() - file_mtime).days} days)")
                except Exception as file_error:
                    logger.warning(f"Error deleting old report {report_file.name}: {file_error}")
        
        if deleted_count > 0:
            logger.info(f"Cleaned up {deleted_count} old report(s) from {reports_dir}")
        
        return deleted_count
    except Exception as e:
        logger.warning(f"Error during report cleanup: {e}")
        return 0


class TaskResult:
    """Result of a task execution"""
    def __init__(self, success: bool, message: str, data: Any = None, error: str = None):
        self.success = success
        self.message = message
        self.data = data
        self.error = error
    
    def to_dict(self):
        return {
            'success': self.success,
            'message': self.message,
            'data': self.data,
            'error': self.error
        }


class BaseTask(ABC):
    """Base class for all tasks that Lea can perform"""
    
    def __init__(self, name: str, description: str, requires_confirmation: bool = False):
        self.name = name
        self.description = description
        self.requires_confirmation = requires_confirmation
        self.allowed = True  # Can be disabled via config
    
    @abstractmethod
    def execute(self, **kwargs) -> TaskResult:
        """Execute the task. Must be implemented by subclasses."""
        pass
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        """Validate task parameters. Override if needed."""
        return True, ""
    
    def get_required_params(self) -> List[str]:
        """Return list of required parameter names. Override if needed."""
        return []


# =====================================================
# FILE OPERATIONS TASKS
# =====================================================

class FileCopyTask(BaseTask):
    """Copy files from source to destination"""
    
    def __init__(self):
        super().__init__(
            name="file_copy",
            description="Copy a file from source to destination",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["source", "destination"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        source = kwargs.get("source")
        destination = kwargs.get("destination")
        
        if not source:
            return False, "Source path is required"
        if not destination:
            return False, "Destination path is required"
        if not os.path.exists(source):
            return False, f"Source file does not exist: {source}"
        
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            source = kwargs.get("source")
            destination = kwargs.get("destination")
            
            # Validate
            valid, msg = self.validate_params(**kwargs)
            if not valid:
                return TaskResult(False, msg, error=msg)
            
            # Ensure destination directory exists
            dest_path = Path(destination)
            dest_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Copy file
            shutil.copy2(source, destination)
            
            return TaskResult(
                True,
                f"File copied successfully from {source} to {destination}",
                data={"source": source, "destination": destination}
            )
        except Exception as e:
            logger.error(f"FileCopyTask error: {e}")
            return TaskResult(False, f"Failed to copy file: {str(e)}", error=str(e))


class FileMoveTask(BaseTask):
    """Move files from source to destination"""
    
    def __init__(self):
        super().__init__(
            name="file_move",
            description="Move a file from source to destination",
            requires_confirmation=True
        )
    
    def get_required_params(self) -> List[str]:
        return ["source", "destination"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        source = kwargs.get("source")
        destination = kwargs.get("destination")
        
        if not source:
            return False, "Source path is required"
        if not destination:
            return False, "Destination path is required"
        if not os.path.exists(source):
            return False, f"Source file does not exist: {source}"
        
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            source = kwargs.get("source")
            destination = kwargs.get("destination")
            
            valid, msg = self.validate_params(**kwargs)
            if not valid:
                return TaskResult(False, msg, error=msg)
            
            dest_path = Path(destination)
            dest_path.parent.mkdir(parents=True, exist_ok=True)
            
            shutil.move(source, destination)
            
            return TaskResult(
                True,
                f"File moved successfully from {source} to {destination}",
                data={"source": source, "destination": destination}
            )
        except Exception as e:
            logger.error(f"FileMoveTask error: {e}")
            return TaskResult(False, f"Failed to move file: {str(e)}", error=str(e))


class FileDeleteTask(BaseTask):
    """Delete files or directories"""
    
    def __init__(self):
        super().__init__(
            name="file_delete",
            description="Delete a file or directory",
            requires_confirmation=True
        )
    
    def get_required_params(self) -> List[str]:
        return ["path"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        path = kwargs.get("path")
        
        if not path:
            return False, "Path is required"
        if not os.path.exists(path):
            return False, f"Path does not exist: {path}"
        
        # Safety check - prevent deletion of critical paths
        critical_paths = [
            os.path.expanduser("~/.ssh"),
            os.path.expanduser("~/Documents"),
            "C:\\Windows",
            "/etc",
            "/usr",
            "/bin",
        ]
        
        abs_path = os.path.abspath(path)
        for critical in critical_paths:
            if critical in abs_path:
                return False, f"Cannot delete critical path: {path}"
        
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            path = kwargs.get("path")
            
            valid, msg = self.validate_params(**kwargs)
            if not valid:
                return TaskResult(False, msg, error=msg)
            
            path_obj = Path(path)
            if path_obj.is_file():
                path_obj.unlink()
                action = "deleted"
            elif path_obj.is_dir():
                shutil.rmtree(path)
                action = "deleted directory"
            else:
                return TaskResult(False, f"Path is neither file nor directory: {path}")
            
            return TaskResult(
                True,
                f"Successfully {action}: {path}",
                data={"path": path, "action": action}
            )
        except Exception as e:
            logger.error(f"FileDeleteTask error: {e}")
            return TaskResult(False, f"Failed to delete: {str(e)}", error=str(e))


class FileReadTask(BaseTask):
    """Read file contents"""
    
    def __init__(self):
        super().__init__(
            name="file_read",
            description="Read contents of a text file",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["path"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        path = kwargs.get("path")
        
        if not path:
            return False, "Path is required"
        if not os.path.exists(path):
            return False, f"File does not exist: {path}"
        if not os.path.isfile(path):
            return False, f"Path is not a file: {path}"
        
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            path = kwargs.get("path")
            encoding = kwargs.get("encoding", "utf-8")
            max_size = kwargs.get("max_size", 10 * 1024 * 1024)  # 10MB default
            
            valid, msg = self.validate_params(**kwargs)
            if not valid:
                return TaskResult(False, msg, error=msg)
            
            # Check file size
            file_size = os.path.getsize(path)
            if file_size > max_size:
                return TaskResult(
                    False,
                    f"File too large ({file_size} bytes). Max size: {max_size} bytes",
                    error="File too large"
                )
            
            with open(path, 'r', encoding=encoding, errors='ignore') as f:
                content = f.read()
            
            return TaskResult(
                True,
                f"Successfully read file: {path}",
                data={"path": path, "content": content, "size": file_size}
            )
        except Exception as e:
            logger.error(f"FileReadTask error: {e}")
            return TaskResult(False, f"Failed to read file: {str(e)}", error=str(e))


class FileWriteTask(BaseTask):
    """Write content to a file"""
    
    def __init__(self):
        super().__init__(
            name="file_write",
            description="Write content to a file",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["path", "content"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        path = kwargs.get("path")
        
        if not path:
            return False, "Path is required"
        
        # Safety check
        critical_paths = ["C:\\Windows", "/etc", "/usr", "/bin"]
        abs_path = os.path.abspath(path)
        for critical in critical_paths:
            if critical in abs_path:
                return False, f"Cannot write to critical path: {path}"
        
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            path = kwargs.get("path")
            content = kwargs.get("content", "")
            encoding = kwargs.get("encoding", "utf-8")
            mode = kwargs.get("mode", "w")  # 'w' for write, 'a' for append
            
            valid, msg = self.validate_params(**kwargs)
            if not valid:
                return TaskResult(False, msg, error=msg)
            
            # Ensure directory exists
            path_obj = Path(path)
            path_obj.parent.mkdir(parents=True, exist_ok=True)
            
            with open(path, mode, encoding=encoding) as f:
                f.write(str(content))
            
            return TaskResult(
                True,
                f"Successfully wrote to file: {path}",
                data={"path": path, "mode": mode, "size": len(str(content))}
            )
        except Exception as e:
            logger.error(f"FileWriteTask error: {e}")
            return TaskResult(False, f"Failed to write file: {str(e)}", error=str(e))


# =====================================================
# DIRECTORY OPERATIONS TASKS
# =====================================================

class DirectoryCreateTask(BaseTask):
    """Create directories"""
    
    def __init__(self):
        super().__init__(
            name="directory_create",
            description="Create a directory (and parent directories if needed)",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["path"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            path = kwargs.get("path")
            if not path:
                return TaskResult(False, "Path is required", error="Missing path")
            
            Path(path).mkdir(parents=True, exist_ok=True)
            
            return TaskResult(
                True,
                f"Directory created (or already exists): {path}",
                data={"path": path}
            )
        except Exception as e:
            logger.error(f"DirectoryCreateTask error: {e}")
            return TaskResult(False, f"Failed to create directory: {str(e)}", error=str(e))


class DirectoryListTask(BaseTask):
    """List directory contents"""
    
    def __init__(self):
        super().__init__(
            name="directory_list",
            description="List files and directories in a path",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["path"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            path = kwargs.get("path", ".")
            if not os.path.exists(path):
                return TaskResult(False, f"Path does not exist: {path}", error="Path not found")
            
            items = []
            for item in os.listdir(path):
                item_path = os.path.join(path, item)
                items.append({
                    "name": item,
                    "path": item_path,
                    "type": "directory" if os.path.isdir(item_path) else "file",
                    "size": os.path.getsize(item_path) if os.path.isfile(item_path) else None
                })
            
            return TaskResult(
                True,
                f"Listed {len(items)} items in {path}",
                data={"path": path, "items": items, "count": len(items)}
            )
        except Exception as e:
            logger.error(f"DirectoryListTask error: {e}")
            return TaskResult(False, f"Failed to list directory: {str(e)}", error=str(e))


# =====================================================
# TEXT PROCESSING TASKS
# =====================================================

class TextReplaceTask(BaseTask):
    """Replace text in a file"""
    
    def __init__(self):
        super().__init__(
            name="text_replace",
            description="Find and replace text in a file",
            requires_confirmation=True
        )
    
    def get_required_params(self) -> List[str]:
        return ["path", "old_text", "new_text"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            path = kwargs.get("path")
            old_text = kwargs.get("old_text")
            new_text = kwargs.get("new_text", "")
            encoding = kwargs.get("encoding", "utf-8")
            
            if not path or old_text is None:
                return TaskResult(False, "Path and old_text are required", error="Missing parameters")
            
            if not os.path.exists(path):
                return TaskResult(False, f"File does not exist: {path}", error="File not found")
            
            # Read file
            with open(path, 'r', encoding=encoding) as f:
                content = f.read()
            
            # Replace
            if old_text not in content:
                return TaskResult(False, f"Text not found in file: {old_text}", error="Text not found")
            
            new_content = content.replace(old_text, new_text)
            replacements = content.count(old_text)
            
            # Write back
            with open(path, 'w', encoding=encoding) as f:
                f.write(new_content)
            
            return TaskResult(
                True,
                f"Replaced '{old_text}' with '{new_text}' ({replacements} occurrences) in {path}",
                data={"path": path, "replacements": replacements}
            )
        except Exception as e:
            logger.error(f"TextReplaceTask error: {e}")
            return TaskResult(False, f"Failed to replace text: {str(e)}", error=str(e))


# =====================================================
# SYSTEM COMMAND TASK (with safety)
# =====================================================

class SystemCommandTask(BaseTask):
    """Execute system commands (safely)"""
    
    def __init__(self):
        super().__init__(
            name="system_command",
            description="Execute a system command (whitelist only)",
            requires_confirmation=True
        )
        
        # Whitelist of allowed commands/patterns
        self.allowed_commands = [
            "python", "python3",
            "pip", "pip3",
            "git",
            "echo",
            "dir", "ls",
            "cd",
        ]
        
        # Blocked commands
        self.blocked_commands = [
            "rm", "del", "format", "fdisk",
            "sudo", "su",
            "chmod 777", "chmod +x",
        ]
        
        # Blocked paths - never allow commands targeting system directories
        self.blocked_paths = [
            "C:\\Windows",
            "C:\\Program Files",
            "C:\\Program Files (x86)",
            "C:\\ProgramData",
            "C:\\System32",
            "C:\\SysWOW64",
            "/Windows",
            "/Program Files",
            "/System32",
            "/SysWOW64",
        ]
    
    def get_required_params(self) -> List[str]:
        return ["command"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        command = kwargs.get("command", "")
        
        if not command:
            return False, "Command is required"
        
        command_lower = command.lower()
        
        # Check for blocked commands
        for blocked in self.blocked_commands:
            if blocked in command_lower:
                return False, f"Blocked command: {blocked}"
        
        # Check for blocked paths - never allow commands targeting system directories
        for blocked_path in self.blocked_paths:
            if blocked_path.lower() in command_lower:
                return False, f"Blocked path: {blocked_path} - system directories are protected"
        
        # Additional safety: block any command that looks like it's trying to modify system files
        dangerous_patterns = [
            "format", "fdisk", "diskpart", "reg delete", "reg add",
            "bcdedit", "bootcfg", "sfc /scannow",  # System modification commands
        ]
        for pattern in dangerous_patterns:
            if pattern.lower() in command_lower:
                return False, f"Blocked dangerous pattern: {pattern}"
        
        # For safety, only allow whitelisted commands or explicit approval
        approved = kwargs.get("approved", False)
        if not approved:
            # Check if command starts with allowed command
            first_word = command.split()[0] if command.split() else ""
            if first_word.lower() not in [cmd.lower() for cmd in self.allowed_commands]:
                return False, f"Command not in whitelist. First word: {first_word}. Set approved=True to allow."
        
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            command = kwargs.get("command")
            timeout = kwargs.get("timeout", 30)
            shell = kwargs.get("shell", True)
            
            valid, msg = self.validate_params(**kwargs)
            if not valid:
                return TaskResult(False, msg, error=msg)
            
            # Execute command
            result = subprocess.run(
                command,
                shell=shell,
                capture_output=True,
                text=True,
                timeout=timeout
            )
            
            output = result.stdout + result.stderr
            
            if result.returncode == 0:
                return TaskResult(
                    True,
                    f"Command executed successfully: {command}",
                    data={"command": command, "output": output, "returncode": result.returncode}
                )
            else:
                return TaskResult(
                    False,
                    f"Command failed with return code {result.returncode}: {command}",
                    data={"command": command, "output": output, "returncode": result.returncode},
                    error=output
                )
        except subprocess.TimeoutExpired:
            return TaskResult(False, f"Command timed out after {timeout} seconds", error="Timeout")
        except Exception as e:
            logger.error(f"SystemCommandTask error: {e}")
            return TaskResult(False, f"Failed to execute command: {str(e)}", error=str(e))


# =====================================================
# SCREEN AUTOMATION TASKS
# =====================================================

class ScreenshotTask(BaseTask):
    """Take a screenshot of the screen"""
    
    def __init__(self):
        super().__init__(
            name="screenshot",
            description="Take a screenshot of the entire screen or a region",
            requires_confirmation=False
        )
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import pyautogui
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            # Optional parameters
            region = kwargs.get("region")  # (x, y, width, height)
            save_path = kwargs.get("save_path")
            
            if region and len(region) == 4:
                screenshot = pyautogui.screenshot(region=region)
            else:
                screenshot = pyautogui.screenshot()
            
            if save_path:
                screenshot.save(save_path)
                return TaskResult(
                    True,
                    f"Screenshot saved to: {save_path}",
                    data={"path": save_path, "size": screenshot.size}
                )
            else:
                # Save to temp location
                import tempfile
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                screenshot.save(temp_file.name)
                return TaskResult(
                    True,
                    f"Screenshot taken and saved to: {temp_file.name}",
                    data={"path": temp_file.name, "size": screenshot.size}
                )
        except Exception as e:
            logger.error(f"ScreenshotTask error: {e}")
            return TaskResult(False, f"Failed to take screenshot: {str(e)}", error=str(e))


class ClickTask(BaseTask):
    """Click at a specific location on screen"""
    
    def __init__(self):
        super().__init__(
            name="click",
            description="Click at specific coordinates on the screen",
            requires_confirmation=True  # Require confirmation for safety
        )
    
    def get_required_params(self) -> List[str]:
        return ["x", "y"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        x = kwargs.get("x")
        y = kwargs.get("y")
        
        if x is None or y is None:
            return False, "x and y coordinates are required"
        
        try:
            x, y = int(x), int(y)
            # Get screen size to validate
            try:
                import pyautogui
                screen_width, screen_height = pyautogui.size()
                if x < 0 or x >= screen_width or y < 0 or y >= screen_height:
                    return False, f"Coordinates ({x}, {y}) are outside screen bounds ({screen_width}x{screen_height})"
            except ImportError:
                pass  # Will be caught in execute
        except (ValueError, TypeError):
            return False, "x and y must be valid integers"
        
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            x = int(kwargs.get("x"))
            y = int(kwargs.get("y"))
            button = kwargs.get("button", "left")  # left, right, middle
            clicks = kwargs.get("clicks", 1)
            interval = kwargs.get("interval", 0.0)
            
            pyautogui.click(x, y, button=button, clicks=clicks, interval=interval)
            
            return TaskResult(
                True,
                f"Clicked at ({x}, {y}) with {button} button ({clicks} clicks)",
                data={"x": x, "y": y, "button": button, "clicks": clicks}
            )
        except Exception as e:
            logger.error(f"ClickTask error: {e}")
            return TaskResult(False, f"Failed to click: {str(e)}", error=str(e))


class TypeTask(BaseTask):
    """Type text at the current cursor position"""
    
    def __init__(self):
        super().__init__(
            name="type",
            description="Type text at the current cursor position",
            requires_confirmation=True  # Require confirmation for safety
        )
    
    def get_required_params(self) -> List[str]:
        return ["text"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        text = kwargs.get("text")
        if not text:
            return False, "Text to type is required"
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            text = str(kwargs.get("text", ""))
            interval = kwargs.get("interval", 0.0)  # Delay between keystrokes
            
            pyautogui.write(text, interval=interval)
            
            return TaskResult(
                True,
                f"Typed text: {text[:50]}{'...' if len(text) > 50 else ''}",
                data={"text_length": len(text), "interval": interval}
            )
        except Exception as e:
            logger.error(f"TypeTask error: {e}")
            return TaskResult(False, f"Failed to type text: {str(e)}", error=str(e))


class KeyPressTask(BaseTask):
    """Press a key or key combination"""
    
    def __init__(self):
        super().__init__(
            name="key_press",
            description="Press a key or key combination (e.g., 'enter', 'ctrl+c')",
            requires_confirmation=True
        )
    
    def get_required_params(self) -> List[str]:
        return ["key"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            key = kwargs.get("key", "")
            presses = kwargs.get("presses", 1)
            interval = kwargs.get("interval", 0.0)
            
            if not key:
                return TaskResult(False, "Key is required", error="Missing key")
            
            pyautogui.press(key, presses=presses, interval=interval)
            
            return TaskResult(
                True,
                f"Pressed key: {key} ({presses} times)",
                data={"key": key, "presses": presses}
            )
        except Exception as e:
            logger.error(f"KeyPressTask error: {e}")
            return TaskResult(False, f"Failed to press key: {str(e)}", error=str(e))


class HotkeyTask(BaseTask):
    """Press a hotkey combination (e.g., Ctrl+C, Alt+Tab)"""
    
    def __init__(self):
        super().__init__(
            name="hotkey",
            description="Press a hotkey combination (e.g., 'ctrl+c', 'alt+tab')",
            requires_confirmation=True
        )
    
    def get_required_params(self) -> List[str]:
        return ["keys"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            keys = kwargs.get("keys", "")
            if not keys:
                return TaskResult(False, "Keys are required (e.g., 'ctrl+c')", error="Missing keys")
            
            # Parse keys (can be string like "ctrl+c" or list like ["ctrl", "c"])
            if isinstance(keys, str):
                key_list = keys.lower().split("+")
            else:
                key_list = [str(k).lower() for k in keys]
            
            pyautogui.hotkey(*key_list)
            
            return TaskResult(
                True,
                f"Pressed hotkey: {'+'.join(key_list)}",
                data={"keys": key_list}
            )
        except Exception as e:
            logger.error(f"HotkeyTask error: {e}")
            return TaskResult(False, f"Failed to press hotkey: {str(e)}", error=str(e))


class FindImageTask(BaseTask):
    """Find an image on the screen"""
    
    def __init__(self):
        super().__init__(
            name="find_image",
            description="Find an image on the screen and return its location",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["image_path"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            image_path = kwargs.get("image_path")
            if not image_path or not os.path.exists(image_path):
                return TaskResult(False, f"Image file not found: {image_path}", error="File not found")
            
            confidence = kwargs.get("confidence", 0.8)
            region = kwargs.get("region")  # Optional region to search
            
            if region and len(region) == 4:
                location = pyautogui.locateOnScreen(image_path, confidence=confidence, region=region)
            else:
                location = pyautogui.locateOnScreen(image_path, confidence=confidence)
            
            if location:
                center = pyautogui.center(location)
                return TaskResult(
                    True,
                    f"Found image at {center} (confidence: {confidence})",
                    data={
                        "x": center.x,
                        "y": center.y,
                        "left": location.left,
                        "top": location.top,
                        "width": location.width,
                        "height": location.height,
                        "confidence": confidence
                    }
                )
            else:
                return TaskResult(
                    False,
                    f"Image not found on screen (confidence: {confidence})",
                    error="Image not found"
                )
        except Exception as e:
            logger.error(f"FindImageTask error: {e}")
            return TaskResult(False, f"Failed to find image: {str(e)}", error=str(e))


class ScrollTask(BaseTask):
    """Scroll the screen"""
    
    def __init__(self):
        super().__init__(
            name="scroll",
            description="Scroll up or down on the screen",
            requires_confirmation=False
        )
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            clicks = kwargs.get("clicks", 3)  # Positive = scroll up, negative = scroll down
            x = kwargs.get("x")  # Optional x coordinate
            y = kwargs.get("y")  # Optional y coordinate
            
            if x is not None and y is not None:
                pyautogui.scroll(clicks, x=x, y=y)
            else:
                pyautogui.scroll(clicks)
            
            direction = "up" if clicks > 0 else "down"
            return TaskResult(
                True,
                f"Scrolled {direction} {abs(clicks)} clicks",
                data={"clicks": clicks, "direction": direction}
            )
        except Exception as e:
            logger.error(f"ScrollTask error: {e}")
            return TaskResult(False, f"Failed to scroll: {str(e)}", error=str(e))


class MoveMouseTask(BaseTask):
    """Move mouse to a specific location"""
    
    def __init__(self):
        super().__init__(
            name="move_mouse",
            description="Move mouse cursor to specific coordinates",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["x", "y"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            x = int(kwargs.get("x"))
            y = int(kwargs.get("y"))
            duration = kwargs.get("duration", 0.5)  # Time to move in seconds
            
            pyautogui.moveTo(x, y, duration=duration)
            
            return TaskResult(
                True,
                f"Moved mouse to ({x}, {y})",
                data={"x": x, "y": y, "duration": duration}
            )
        except Exception as e:
            logger.error(f"MoveMouseTask error: {e}")
            return TaskResult(False, f"Failed to move mouse: {str(e)}", error=str(e))


class GetScreenSizeTask(BaseTask):
    """Get the screen size"""
    
    def __init__(self):
        super().__init__(
            name="get_screen_size",
            description="Get the screen resolution (width x height)",
            requires_confirmation=False
        )
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            try:
                import pyautogui
            except ImportError:
                return TaskResult(
                    False,
                    "pyautogui not installed. Install with: pip install pyautogui",
                    error="Module not found"
                )
            
            width, height = pyautogui.size()
            
            return TaskResult(
                True,
                f"Screen size: {width}x{height}",
                data={"width": width, "height": height}
            )
        except Exception as e:
            logger.error(f"GetScreenSizeTask error: {e}")
            return TaskResult(False, f"Failed to get screen size: {str(e)}", error=str(e))


# =====================================================
# EMAIL TASKS (Microsoft Graph / Outlook)
# =====================================================

class OutlookEmailCheckTask(BaseTask):
    """Check Outlook inbox and generate email report"""
    
    def __init__(self):
        super().__init__(
            name="outlook_email_check",
            description="Check Outlook inbox and generate email report",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return []
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                # Fallback: try outlook_email module
                try:
                    from outlook_email import get_email_client
                    client = get_email_client()
                except ImportError:
                    return TaskResult(
                        False,
                        "Outlook integration not available. Please ensure outlook_integration.py or outlook_email.py exists.",
                        error="Module not found"
                    )
            
            if not client:
                return TaskResult(
                    False,
                    "Outlook client not initialized. Check your .env file for OUTLOOK_CLIENT_ID and OUTLOOK_TENANT_ID. OUTLOOK_CLIENT_SECRET is optional (not needed for manual login).",
                    error="Client not initialized"
                )
            
            # Authenticate if needed
            try:
                if not client.is_authenticated():
                    logger.info("🔐 Authenticating with Outlook...")
                    logger.info("A browser window will open for you to sign in with your Microsoft account.")
                    auth_result = client.authenticate()
                    if not auth_result:
                        # Get more detailed error from logs - the authenticate method already logged details
                        return TaskResult(
                            False,
                            "Authentication failed. Check the logs for detailed error information.\n\nCommon issues:\n1. Browser login window didn't open or wasn't completed\n2. Permissions not granted in Azure AD app registration\n3. User account not found in the specified tenant\n\nCheck lea_crash.log for detailed error messages.",
                            error="Authentication failed"
                        )
                else:
                    logger.info("✅ Already authenticated with Outlook")
            except AttributeError:
                # Client might not have is_authenticated method (old interface)
                logger.info("🔐 Authenticating with Outlook...")
                logger.info("A browser window will open for you to sign in with your Microsoft account.")
                auth_result = client.authenticate()
                if not auth_result:
                    return TaskResult(
                        False,
                        "Authentication failed. Check the logs for detailed error information.\n\nCommon issues:\n1. Browser login window didn't open or wasn't completed\n2. Permissions not granted in Azure AD app registration\n3. User account not found in the specified tenant\n\nCheck lea_crash.log for detailed error messages.",
                        error="Authentication failed"
                    )
            
            # Get parameters
            max_results = kwargs.get("max_results", 1000)  # Increased default for full inbox
            timeframe_days = kwargs.get("timeframe_days", None)  # e.g., 365 for last year
            include_folders = kwargs.get("include_folders", False)  # Default: inbox only
            generate_analysis = kwargs.get("generate_analysis", False)
            schedule_regular = kwargs.get("schedule_regular", False)
            
            # Calculate date filter if timeframe specified
            date_filter = None
            if timeframe_days:
                from datetime import timedelta
                cutoff_date = datetime.now() - timedelta(days=int(timeframe_days))
                date_filter = cutoff_date.isoformat()
                logger.info(f"Filtering emails from last {timeframe_days} days (since {cutoff_date.date()})")
            
            # Get emails (inbox only if include_folders is False)
            # Try to pass date filter if client supports it, otherwise filter after retrieval
            try:
                if include_folders:
                    # Get from all folders - this might need folder-specific API calls
                    # Try with date filter first
                    try:
                        emails = client.get_emails(max_results=max_results, since=date_filter)
                    except TypeError:
                        # Client doesn't support 'since' parameter, get all and filter later
                        emails = client.get_emails(max_results=max_results)
                else:
                    # Get from inbox only (default behavior)
                    try:
                        emails = client.get_emails(max_results=max_results, since=date_filter)
                    except TypeError:
                        # Client doesn't support 'since' parameter, get all and filter later
                        emails = client.get_emails(max_results=max_results)
            except Exception as get_emails_error:
                logger.warning(f"Error getting emails with parameters: {get_emails_error}, trying without date filter")
                # Fallback: get emails without date filter
                emails = client.get_emails(max_results=max_results)
            
            if not emails:
                return TaskResult(
                    True,
                    "No emails found in inbox" + (f" from the last {timeframe_days} days" if timeframe_days else ""),
                    data={"emails": [], "count": 0}
                )
            
            # Filter by date if needed (in case client doesn't support date filtering)
            if date_filter:
                filtered_emails = []
                cutoff_datetime = datetime.fromisoformat(date_filter.replace('Z', '+00:00'))
                for email in emails:
                    email_date_str = email.get('receivedDateTime') or email.get('sentDateTime') or email.get('date', '')
                    if email_date_str:
                        try:
                            # Parse email date
                            if 'T' in email_date_str:
                                email_date = datetime.fromisoformat(email_date_str.replace('Z', '+00:00'))
                            else:
                                email_date = datetime.strptime(email_date_str, '%Y-%m-%d')
                            
                            # Compare (handle timezone-aware dates)
                            if hasattr(email_date, 'replace'):
                                email_date_naive = email_date.replace(tzinfo=None) if email_date.tzinfo else email_date
                                cutoff_naive = cutoff_datetime.replace(tzinfo=None) if cutoff_datetime.tzinfo else cutoff_datetime
                                if email_date_naive >= cutoff_naive:
                                    filtered_emails.append(email)
                        except Exception as date_error:
                            logger.warning(f"Error parsing email date {email_date_str}: {date_error}")
                            # Include email if date parsing fails (better to include than exclude)
                            filtered_emails.append(email)
                    else:
                        # Include if no date (better to include than exclude)
                        filtered_emails.append(email)
                emails = filtered_emails
                logger.info(f"Filtered to {len(emails)} emails after date filtering")
            
            # Process emails to create enhanced report with requested columns
            import re
            import html
            from datetime import datetime as dt
            
            processed_emails = []
            for email in emails:
                # Format date properly
                date_sent = email.get('sentDateTime') or email.get('receivedDateTime') or email.get('date', '')
                if date_sent:
                    try:
                        # Try to parse ISO format date and format it nicely
                        if isinstance(date_sent, str):
                            # Parse ISO format (e.g., "2024-11-23T10:30:00Z" or "2024-11-23T10:30:00+00:00")
                            if 'T' in date_sent:
                                # Remove timezone info for parsing
                                date_str = date_sent.replace('Z', '+00:00')
                                try:
                                    date_obj = datetime.fromisoformat(date_str)
                                except ValueError:
                                    # Try without timezone
                                    date_str_clean = date_str.split('+')[0].split('-')[0] if '+' in date_str else date_str
                                    date_obj = datetime.fromisoformat(date_str_clean)
                                date_sent = date_obj.strftime('%Y-%m-%d %H:%M:%S')
                            # else: already formatted, use as-is
                    except Exception as date_error:
                        # If parsing fails, use as-is
                        logger.debug(f"Date parsing failed for {date_sent}: {date_error}")
                        pass
                else:
                    date_sent = 'N/A'
                
                processed = {
                    'Date Sent': date_sent,
                    'Sender': email.get('from', {}).get('emailAddress', {}).get('address', '') if isinstance(email.get('from'), dict) else str(email.get('from', '')),
                    'Subject': email.get('subject', ''),
                    'Synopsis': '',  # Will be generated from body
                    'Mentions Andrea/Dre': '',
                    'Tasks Requested': '',
                    'Has Attachments': 'No'
                }
                
                # Check for attachments
                has_attachments = email.get('hasAttachments', False)
                if has_attachments:
                    processed['Has Attachments'] = 'Yes'
                
                # Get email body for synopsis and task detection
                body = email.get('body', {}).get('content', '') if isinstance(email.get('body'), dict) else str(email.get('body', ''))
                body_preview = email.get('bodyPreview', '')
                full_body = body or body_preview or ''
                
                # Create synopsis (2-3 sentence summary from body)
                if full_body:
                    # Remove HTML tags and decode HTML entities
                    text_body = re.sub(r'<[^>]+>', '', full_body)
                    # Decode HTML entities (like &nbsp;, &amp;, etc.)
                    text_body = html.unescape(text_body)
                    # Replace multiple whitespace with single space
                    text_body = re.sub(r'\s+', ' ', text_body).strip()
                    
                    # Create a proper summary (2-3 sentences)
                    # Strategy: Find the most important sentences (usually first paragraph or sentences with key words)
                    sentences = re.split(r'[.!?]+', text_body)
                    # Filter out empty sentences and very short ones
                    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
                    
                    if len(sentences) > 0:
                        # Prioritize sentences that mention tasks, actions, or key information
                        priority_sentences = []
                        regular_sentences = []
                        
                        for sentence in sentences:
                            sentence_lower = sentence.lower()
                            # Prioritize sentences with action words, questions, or important keywords
                            if any(word in sentence_lower for word in ['please', 'need', 'can you', 'could you', 'would you', 
                                                                      'task', 'action', 'do', 'complete', 'send', 'provide', 
                                                                      'request', 'required', 'important', 'urgent', 'deadline',
                                                                      'andrea', 'dre', 'meeting', 'call', 'review', 'update']):
                                priority_sentences.append(sentence)
                            else:
                                regular_sentences.append(sentence)
                        
                        # Build synopsis: prefer priority sentences, fall back to first sentences
                        synopsis_parts = []
                        if priority_sentences:
                            # Take up to 2 priority sentences
                            synopsis_parts.extend(priority_sentences[:2])
                            # If we need a third sentence and have regular sentences, add one
                            if len(synopsis_parts) < 3 and regular_sentences:
                                synopsis_parts.append(regular_sentences[0])
                        else:
                            # No priority sentences, use first 2-3 sentences
                            synopsis_parts.extend(sentences[:3])
                        
                        # Join sentences and ensure proper length (2-3 sentences, max ~300 chars)
                        synopsis = '. '.join(synopsis_parts[:3]).strip()
                        if synopsis and not synopsis.endswith(('.', '!', '?')):
                            synopsis += '.'
                        
                        # Truncate if too long (shouldn't happen with 3 sentences, but safety check)
                        if len(synopsis) > 400:
                            # Find last complete sentence within limit
                            last_period = synopsis[:400].rfind('.')
                            if last_period > 100:  # Ensure we have at least one sentence
                                synopsis = synopsis[:last_period+1]
                            else:
                                synopsis = synopsis[:300].strip() + '...'
                        
                        processed['Synopsis'] = synopsis if synopsis else text_body[:200].strip() + '...'
                    else:
                        # Fallback: just use first 200 chars
                        processed['Synopsis'] = text_body[:200].strip() + '...'
                
                # Check for mentions of "andrea" or "dre" (case-insensitive)
                # Clean the body first (same as synopsis) to remove HTML before checking
                cleaned_body = full_body
                if cleaned_body:
                    # Remove HTML tags
                    cleaned_body = re.sub(r'<[^>]+>', '', cleaned_body)
                    # Decode HTML entities
                    cleaned_body = html.unescape(cleaned_body)
                    # Replace multiple whitespace/newlines with single space
                    cleaned_body = re.sub(r'\s+', ' ', cleaned_body).strip()
                
                body_lower = cleaned_body.lower()
                mentions = []
                if 'andrea' in body_lower or 'dre' in body_lower:
                    processed['Mentions Andrea/Dre'] = 'Yes'
                    # Try to extract task-related sentences from cleaned text
                    sentences = re.split(r'[.!?]+', cleaned_body)
                    task_sentences = []
                    for sentence in sentences:
                        sentence = sentence.strip()
                        if not sentence:  # Skip empty sentences
                            continue
                        sentence_lower = sentence.lower()
                        if ('andrea' in sentence_lower or 'dre' in sentence_lower) and \
                           any(word in sentence_lower for word in ['please', 'need', 'can you', 'could you', 'would you', 'task', 'action', 'do', 'complete']):
                            # Clean up the sentence one more time to remove any remaining artifacts
                            clean_sentence = re.sub(r'\s+', ' ', sentence).strip()
                            if clean_sentence and len(clean_sentence) > 10:  # Only add meaningful sentences
                                task_sentences.append(clean_sentence)
                    if task_sentences:
                        processed['Tasks Requested'] = ' | '.join(task_sentences[:3])  # Limit to 3 task sentences
                        # Also include tasks in synopsis if not already there
                        if processed['Synopsis'] and not any(task in processed['Synopsis'] for task in task_sentences[:1]):
                            # Add the first task sentence to synopsis if it's not already included
                            first_task = task_sentences[0]
                            if len(processed['Synopsis']) + len(first_task) < 400:
                                processed['Synopsis'] += ' ' + first_task + '.'
                else:
                    processed['Mentions Andrea/Dre'] = 'No'
                
                processed_emails.append(processed)
            
            # Generate analysis summary if requested
            analysis_summary = None
            if generate_analysis:
                # Count unread
                unread_count = sum(1 for e in emails if not e.get('isRead', True))
                # Count flagged
                flagged_count = sum(1 for e in emails if isinstance(e.get('flag'), dict) and e.get('flag').get('flagStatus') == 'flagged')
                # Top senders
                sender_counts = {}
                for email in emails:
                    sender = email.get('from', {}).get('emailAddress', {}).get('address', '') if isinstance(email.get('from'), dict) else str(email.get('from', ''))
                    if sender:
                        sender_counts[sender] = sender_counts.get(sender, 0) + 1
                top_senders = sorted(sender_counts.items(), key=lambda x: x[1], reverse=True)[:10]
                
                analysis_summary = {
                    'total_emails': len(emails),
                    'unread_count': unread_count,
                    'flagged_count': flagged_count,
                    'top_senders': top_senders
                }
            
            # Generate CSV report (always works - uses only built-in Python csv module)
            reports_dir = REPORTS_DIR
            reports_dir.mkdir(parents=True, exist_ok=True)
            
            # Clean up old reports (older than 10 days)
            cleanup_old_reports(reports_dir, days_old=10)
            
            import csv
            
            report_path = reports_dir / f"email_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
            
            if not processed_emails:
                return TaskResult(
                    True,
                    f"No emails found to generate report.",
                    data={"emails": processed_emails, "count": len(emails), "analysis": analysis_summary}
                )
            
            # Get all column names from first email dict
            fieldnames = list(processed_emails[0].keys())
            
            # Write CSV using built-in csv module (no external dependencies)
            with open(report_path, 'w', newline='', encoding='utf-8') as csvfile:
                writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
                writer.writeheader()
                writer.writerows(processed_emails)
            
            logger.info(f"CSV report generated: {report_path}")
            
            # Build result message
            result_message = f"Email report generated (CSV): {len(emails)} emails found. Report saved to: {report_path}"
            
            if analysis_summary:
                result_message += f"\n\nAnalysis Summary:\n- Total: {analysis_summary['total_emails']} emails\n- Unread: {analysis_summary['unread_count']}\n- Flagged: {analysis_summary['flagged_count']}"
                if analysis_summary['top_senders']:
                    result_message += f"\n- Top sender: {analysis_summary['top_senders'][0][0]} ({analysis_summary['top_senders'][0][1]} emails)"
            
            # Note about scheduling
            if schedule_regular:
                result_message += "\n\nNote: Regular scheduling is not yet implemented. You can manually run this task periodically, or we can implement scheduling in a future update."
            
            return TaskResult(
                True,
                result_message,
                data={
                    "emails": processed_emails,
                    "count": len(emails),
                    "report_path": str(report_path),
                    "report_type": "csv",
                    "analysis": analysis_summary
                }
            )
        except Exception as e:
            # Use logger if available, otherwise use print as fallback
            try:
                logger.error(f"OutlookEmailCheckTask error: {e}")
            except NameError:
                # Fallback if logger is not available
                print(f"OutlookEmailCheckTask error: {e}")
                import traceback
                traceback.print_exc()
            return TaskResult(False, f"Failed to check Outlook emails: {str(e)}", error=str(e))


class OutlookEmailDraftTask(BaseTask):
    """Create a draft email in Outlook"""
    
    def __init__(self):
        super().__init__(
            name="outlook_email_draft",
            description="Create a draft email in Outlook",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["subject", "body"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                # Fallback: try outlook_email module
                try:
                    from outlook_email import get_email_client
                    client = get_email_client()
                except ImportError:
                    return TaskResult(
                        False,
                        "Outlook integration not available. Please ensure outlook_integration.py or outlook_email.py exists.",
                        error="Module not found"
                    )
            
            if not client:
                return TaskResult(
                    False,
                    "Outlook client not initialized. Check your .env file for Outlook credentials.",
                    error="Client not initialized"
                )
            
            # Authenticate if needed
            try:
                if not client.is_authenticated():
                    logger.info("🔐 Authenticating with Outlook...")
                    logger.info("A browser window will open for you to sign in with your Microsoft account.")
                    if not client.authenticate():
                        return TaskResult(
                            False,
                            "Authentication failed. Please:\n1. Check OUTLOOK_CLIENT_ID and OUTLOOK_TENANT_ID in .env\n2. Complete the browser login when prompted\n3. Grant the requested permissions",
                            error="Authentication failed"
                        )
            except AttributeError:
                # Client might not have is_authenticated method (old interface)
                logger.info("🔐 Authenticating with Outlook...")
                logger.info("A browser window will open for you to sign in with your Microsoft account.")
                if not client.authenticate():
                    return TaskResult(
                        False,
                        "Authentication failed. Please:\n1. Check OUTLOOK_CLIENT_ID and OUTLOOK_TENANT_ID in .env\n2. Complete the browser login when prompted\n3. Grant the requested permissions",
                        error="Authentication failed"
                    )
            
            subject = kwargs.get("subject", "").strip()
            body = kwargs.get("body", "").strip()
            to = kwargs.get("to", "").strip()
            cc = kwargs.get("cc", "").strip()
            bcc = kwargs.get("bcc", "").strip()
            
            # Check for missing required parameters and provide helpful error
            missing_params = []
            if not subject:
                missing_params.append("subject")
            if not body:
                missing_params.append("body")
            
            if missing_params:
                error_msg = f"Missing required parameters: {', '.join(missing_params)}.\n\n"
                error_msg += "Please extract the subject and body from the conversation history, or ask the user for these details.\n"
                error_msg += "If the user mentioned draft details earlier (like 'Subject: ...' or 'Body: ...'), extract those values from previous messages.\n"
                error_msg += "Then call this task again with the extracted parameters."
                return TaskResult(
                    False,
                    error_msg,
                    error=f"Missing required parameters: {', '.join(missing_params)}"
                )
            
            # Create draft
            draft_id = client.create_draft(
                subject=subject,
                body=body,
                to=to,
                cc=cc,
                bcc=bcc
            )
            
            if draft_id:
                return TaskResult(
                    True,
                    f"Draft email created successfully: '{subject}'. Draft ID: {draft_id}",
                    data={"draft_id": draft_id, "subject": subject}
                )
            else:
                return TaskResult(
                    False,
                    "Failed to create draft email",
                    error="Draft creation failed"
                )
        except Exception as e:
            logger.error(f"OutlookEmailDraftTask error: {e}")
            return TaskResult(False, f"Failed to create draft email: {str(e)}", error=str(e))


class OutlookCalendarCheckTask(BaseTask):
    """Check Outlook calendar and generate calendar report"""
    
    def __init__(self):
        super().__init__(
            name="outlook_calendar_check",
            description="Check Outlook calendar and generate calendar report",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return []
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                return TaskResult(
                    False,
                    "Outlook integration not available. Please ensure outlook_integration.py exists.",
                    error="Module not found"
                )
            
            if not client:
                return TaskResult(
                    False,
                    "Outlook client not initialized. Check your .env file for Outlook credentials.",
                    error="Client not initialized"
                )
            
            # Authenticate if needed
            try:
                if not client.is_authenticated():
                    logger.info("🔐 Authenticating with Outlook...")
                    logger.info("A browser window will open for you to sign in with your Microsoft account.")
                    if not client.authenticate():
                        return TaskResult(
                            False,
                            "Authentication failed. Please:\n1. Check OUTLOOK_CLIENT_ID and OUTLOOK_TENANT_ID in .env\n2. Complete the browser login when prompted\n3. Grant the requested permissions",
                            error="Authentication failed"
                        )
            except AttributeError:
                # Client might not have is_authenticated method (old interface)
                logger.info("🔐 Authenticating with Outlook...")
                logger.info("A browser window will open for you to sign in with your Microsoft account.")
                if not client.authenticate():
                    return TaskResult(
                        False,
                        "Authentication failed. Please:\n1. Check OUTLOOK_CLIENT_ID and OUTLOOK_TENANT_ID in .env\n2. Complete the browser login when prompted\n3. Grant the requested permissions",
                        error="Authentication failed"
                    )
            
            # Get calendar events - check if days_ahead parameter is provided
            days_ahead = kwargs.get("days_ahead", 1)  # Default to 1 day (next day only)
            
            # Get calendar events
            events = client.get_calendar_events(days_ahead=days_ahead)
            
            # Filter to only next day if this is an automatic report
            if days_ahead == 1:
                from datetime import datetime, timedelta
                tomorrow = datetime.now() + timedelta(days=1)
                tomorrow_start = tomorrow.replace(hour=0, minute=0, second=0, microsecond=0)
                tomorrow_end = tomorrow.replace(hour=23, minute=59, second=59, microsecond=999999)
                
                # Filter events to only tomorrow
                filtered_events = []
                for event in events:
                    event_start = event.get('start')
                    if event_start:
                        try:
                            # Parse event start time
                            if isinstance(event_start, str):
                                event_dt = datetime.fromisoformat(event_start.replace('Z', '+00:00').replace('+00:00', ''))
                            else:
                                event_dt = event_start
                            
                            # Check if event is tomorrow
                            if tomorrow_start.date() <= event_dt.date() <= tomorrow_end.date():
                                filtered_events.append(event)
                        except Exception as e:
                            logger.debug(f"Error parsing event date {event_start}: {e}")
                            # Include event if date parsing fails (better to include than exclude)
                            filtered_events.append(event)
                    else:
                        # Include if no start time
                        filtered_events.append(event)
                
                events = filtered_events
            
            # Only generate report if there are events
            if not events:
                return TaskResult(
                    True,
                    f"No calendar events found for the next {days_ahead} day(s). No report generated.",
                    data={"events": [], "count": 0}
                )
            
            # Generate CSV report (always works - uses only built-in Python csv module)
            import csv
            from pathlib import Path
            
            # Use the centralized reports directory
            reports_dir = REPORTS_DIR
            reports_dir.mkdir(parents=True, exist_ok=True)
            
            # Clean up old reports (older than 10 days)
            cleanup_old_reports(reports_dir, days_old=10)
            
            if not events:
                return TaskResult(
                    True,
                    f"No calendar events found to generate report.",
                    data={"events": events, "count": len(events)}
                )
            
            # Get all field names from first event
            fieldnames = list(events[0].keys()) if events else []
            
            report_path = reports_dir / f"calendar_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
            
            # Write CSV using built-in csv module (no external dependencies)
            with open(report_path, 'w', newline='', encoding='utf-8') as csvfile:
                writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
                writer.writeheader()
                writer.writerows(events)
            
            logger.info(f"CSV calendar report generated: {report_path}")
            
            return TaskResult(
                True,
                f"Calendar report generated (CSV): {len(events)} events found. Report saved to: {report_path}",
                data={"events": events, "count": len(events), "report_path": str(report_path)}
            )
        except Exception as e:
            logger.error(f"OutlookCalendarCheckTask error: {e}")
            return TaskResult(False, f"Failed to check calendar: {str(e)}", error=str(e))


class OutlookInboxOrganizeTask(BaseTask):
    """Organize Outlook inbox and folders"""
    
    def __init__(self):
        super().__init__(
            name="outlook_inbox_organize",
            description="Organize and clean Outlook inbox and folders",
            requires_confirmation=True  # Require confirmation for organization
        )
    
    def get_required_params(self) -> List[str]:
        return ["action"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                return TaskResult(
                    False,
                    "Outlook integration not available. Please ensure outlook_integration.py exists.",
                    error="Module not found"
                )
            
            if not client:
                return TaskResult(
                    False,
                    "Outlook client not initialized. Check your .env file for Outlook credentials.",
                    error="Client not initialized"
                )
            
            action = kwargs.get("action", "plan")
            folder = kwargs.get("folder", "inbox")
            rules = kwargs.get("rules", {})
            
            if action == "plan":
                # Create organization plan
                plan = client.create_organization_plan(folder=folder, rules=rules)
                return TaskResult(
                    True,
                    f"Organization plan created for {folder}. Review the plan and use action='execute' to apply it.",
                    data={"plan": plan, "folder": folder}
                )
            elif action == "execute":
                # Execute organization
                result = client.organize_inbox(folder=folder, rules=rules)
                return TaskResult(
                    True,
                    f"Inbox organization completed for {folder}",
                    data={"result": result, "folder": folder}
                )
            else:
                return TaskResult(
                    False,
                    f"Invalid action: {action}. Use 'plan' or 'execute'",
                    error="Invalid action"
                )
        except Exception as e:
            logger.error(f"OutlookInboxOrganizeTask error: {e}")
            return TaskResult(False, f"Failed to organize inbox: {str(e)}", error=str(e))


class OutlookUserProfileTask(BaseTask):
    """Get or update Outlook user profile information"""
    
    def __init__(self):
        super().__init__(
            name="outlook_user_profile",
            description="Get or update Outlook user profile information",
            requires_confirmation=False  # Read doesn't need confirmation, update does
        )
    
    def get_required_params(self) -> List[str]:
        return ["action"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                return TaskResult(
                    False,
                    "Outlook integration not available. Please ensure outlook_integration.py exists.",
                    error="Module not found"
                )
            
            if not client:
                return TaskResult(
                    False,
                    "Outlook client not initialized. Check your .env file for Outlook credentials.",
                    error="Client not initialized"
                )
            
            action = kwargs.get("action", "read")
            
            if action == "read":
                # Get profile information
                profile = client.get_user_profile()
                return TaskResult(
                    True,
                    f"User profile retrieved: {profile.get('displayName', 'Unknown')} ({profile.get('mail', 'No email')})",
                    data={"profile": profile}
                )
            elif action == "update":
                # Update profile (requires confirmation)
                updates = kwargs.get("updates", {})
                if not updates:
                    return TaskResult(
                        False,
                        "No updates provided. Use 'updates' parameter with fields to update.",
                        error="Missing updates"
                    )
                
                result = client.update_user_profile(updates=updates)
                return TaskResult(
                    True,
                    f"User profile updated successfully",
                    data={"result": result, "updates": updates}
                )
            else:
                return TaskResult(
                    False,
                    f"Invalid action: {action}. Use 'read' or 'update'",
                    error="Invalid action"
                )
        except Exception as e:
            logger.error(f"OutlookUserProfileTask error: {e}")
            return TaskResult(False, f"Failed to access user profile: {str(e)}", error=str(e))


# Legacy task name for backward compatibility
class EmailCheckTask(BaseTask):
    """Check email inbox for new messages (legacy name - use outlook_email_check instead)"""
    
    def __init__(self):
        super().__init__(
            name="email_check",
            description="Check email inbox for new/unread messages",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return []
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration (preferred)
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                # Fallback: try outlook_email module (legacy)
                try:
                    from outlook_email import get_email_client
                    client = get_email_client()
                except ImportError:
                    return TaskResult(
                        False,
                        "Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv",
                        error="No module named 'outlook_email'"
                    )
            
            max_results = kwargs.get("max_results", 10)
            unread_only = kwargs.get("unread_only", True)
            get_all = kwargs.get("get_all", False)  # New parameter to get all emails
            
            if not client:
                error_msg = """Email client not available. To use Outlook email features:
                
1. Install required packages: pip install msal requests python-dotenv
2. Register an app in Azure AD (https://portal.azure.com):
   - Go to Azure Active Directory > App registrations > New registration
   - Set redirect URI to: http://localhost
   - Add API permissions: Microsoft Graph > Mail.Read, Mail.Send
   - Grant admin consent
3. Copy the Application (client) ID
4. Add to .env file: OUTLOOK_CLIENT_ID=your_client_id_here

For Outlook online, you need to register an Azure AD app to get a CLIENT_ID."""
                return TaskResult(False, error_msg, error="Email client not initialized")
            
            # Authenticate if needed - this will open a browser for login
            if not client.is_authenticated():
                logger.info("🔐 Attempting to authenticate with Outlook...")
                auth_result = client.authenticate(interactive=True)
                if not auth_result:
                    error_msg = """Authentication failed. 
                    
A browser window should have opened for you to sign in to your Microsoft account.
If no browser opened, check:
1. OUTLOOK_CLIENT_ID is set correctly in .env
2. The app is registered in Azure AD with correct permissions
3. Your internet connection is working

Please try again and complete the browser login when prompted."""
                    return TaskResult(False, error_msg, error="Authentication failed")
                logger.info("✅ Authentication successful!")
            
            # Get emails using the OutlookClient interface
            # Note: outlook_integration uses get_emails() method, not check_email()
            if hasattr(client, 'check_email'):
                emails = client.check_email(max_results=max_results, unread_only=unread_only, get_all=get_all)
            elif hasattr(client, 'get_emails'):
                # Use outlook_integration interface
                all_emails = client.get_emails(max_results=max_results)
                if unread_only:
                    emails = [e for e in all_emails if not e.get('isRead', False)]
                else:
                    emails = all_emails
            else:
                return TaskResult(False, "Email client does not support email retrieval", error="Unsupported client interface")
            
            if not emails:
                return TaskResult(True, "No unread emails found" if unread_only else "No emails found", 
                                data={"emails": [], "count": 0})
            
            return TaskResult(True, f"Found {len(emails)} email(s)", 
                            data={"emails": emails, "count": len(emails)})
        except ImportError as ie:
            error_msg = f"Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv\n\nError: {str(ie)}"
            return TaskResult(False, error_msg, error="Import error")
        except Exception as e:
            logger.error(f"EmailCheckTask error: {e}")
            error_msg = f"Failed to check email: {str(e)}\n\nIf this is an authentication error, make sure:\n1. OUTLOOK_CLIENT_ID is set in .env\n2. You completed the browser login\n3. The Azure app has Mail.Read permission"
            return TaskResult(False, error_msg, error=str(e))


class EmailReadTask(BaseTask):
    """Read full content of a specific email"""
    
    def __init__(self):
        super().__init__(
            name="email_read",
            description="Read full content of an email by ID",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["email_id"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration (preferred)
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                # Fallback: try outlook_email module (legacy)
                try:
                    from outlook_email import get_email_client
                    client = get_email_client()
                except ImportError:
                    return TaskResult(
                        False,
                        "Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv",
                        error="Import error"
                    )
            
            email_id = kwargs.get("email_id")
            if not email_id:
                return TaskResult(False, "email_id parameter is required", error="Missing parameter")
            
            if not client:
                return TaskResult(False, "Email client not available", error="Email client not initialized")
            
            if not client.is_authenticated():
                if not client.authenticate():
                    return TaskResult(False, "Email authentication required", error="Authentication failed")
            
            # Try to get email - outlook_integration doesn't have get_email method, so use API directly
            if hasattr(client, 'get_email'):
                email = client.get_email(email_id)
            elif hasattr(client, '_make_request'):
                # Use outlook_integration's _make_request method
                result = client._make_request("GET", f"/me/messages/{email_id}")
                if result:
                    email = {
                        "id": result.get("id"),
                        "subject": result.get("subject", "(No Subject)"),
                        "from": result.get("from", {}).get("emailAddress", {}).get("address", "Unknown"),
                        "body": result.get("body", {}).get("content", ""),
                        "received": result.get("receivedDateTime")
                    }
                else:
                    return TaskResult(False, f"Email with ID {email_id} not found", error="Email not found")
            else:
                return TaskResult(False, "Email client does not support reading individual emails", error="Unsupported client interface")
            
            return TaskResult(True, f"Email retrieved: {email.get('subject', 'No Subject')}", data=email)
        except ImportError as ie:
            return TaskResult(False, f"Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv\n\nError: {str(ie)}", error="Import error")
        except Exception as e:
            logger.error(f"EmailReadTask error: {e}")
            return TaskResult(False, f"Failed to read email: {str(e)}", error=str(e))


class EmailSendTask(BaseTask):
    """Send an email"""
    
    def __init__(self):
        super().__init__(
            name="email_send",
            description="Send an email",
            requires_confirmation=True  # Require confirmation for sending emails
        )
    
    def get_required_params(self) -> List[str]:
        return ["to", "subject", "body"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration (preferred)
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                # Fallback: try outlook_email module (legacy)
                try:
                    from outlook_email import get_email_client
                    client = get_email_client()
                except ImportError:
                    return TaskResult(
                        False,
                        "Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv",
                        error="Import error"
                    )
            
            to = kwargs.get("to")
            subject = kwargs.get("subject")
            body = kwargs.get("body")
            body_type = kwargs.get("body_type", "HTML")
            cc = kwargs.get("cc")
            
            if not to or not subject or not body:
                return TaskResult(False, "to, subject, and body parameters are required", 
                                error="Missing parameters")
            
            if not client:
                return TaskResult(False, "Email client not available", error="Email client not initialized")
            
            if not client.is_authenticated():
                if not client.authenticate():
                    return TaskResult(False, "Email authentication required", error="Authentication failed")
            
            # Use outlook_integration's create_draft method (sends email)
            if hasattr(client, 'send_email'):
                success = client.send_email(to=to, subject=subject, body=body, body_type=body_type, cc=cc)
            elif hasattr(client, 'create_draft'):
                # outlook_integration has create_draft, but we need to send
                # For now, create draft and inform user
                draft_id = client.create_draft(subject=subject, body=body, to=to, cc=cc)
                if draft_id:
                    return TaskResult(True, f"Draft created successfully (ID: {draft_id}). Use Outlook to send it.", 
                                    data={"draft_id": draft_id, "to": to, "subject": subject})
                else:
                    return TaskResult(False, "Failed to create email draft", error="Draft creation failed")
            else:
                return TaskResult(False, "Email client does not support sending emails", error="Unsupported client interface")
            
            if success:
                return TaskResult(True, f"Email sent successfully to {to}", data={"to": to, "subject": subject})
            else:
                return TaskResult(False, "Failed to send email", error="Send failed")
        except ImportError as ie:
            return TaskResult(False, f"Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv\n\nError: {str(ie)}", error="Import error")
        except Exception as e:
            logger.error(f"EmailSendTask error: {e}")
            return TaskResult(False, f"Failed to send email: {str(e)}", error=str(e))


class EmailMarkReadTask(BaseTask):
    """Mark an email as read"""
    
    def __init__(self):
        super().__init__(
            name="email_mark_read",
            description="Mark an email as read",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["email_id"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            # Try to import outlook integration (preferred)
            try:
                from outlook_integration import OutlookClient
                client = OutlookClient()
            except ImportError:
                # Fallback: try outlook_email module (legacy)
                try:
                    from outlook_email import get_email_client
                    client = get_email_client()
                except ImportError:
                    return TaskResult(
                        False,
                        "Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv",
                        error="Import error"
                    )
            
            email_id = kwargs.get("email_id")
            if not email_id:
                return TaskResult(False, "email_id parameter is required", error="Missing parameter")
            
            if not client:
                return TaskResult(False, "Email client not available", error="Email client not initialized")
            
            if not client.is_authenticated():
                if not client.authenticate():
                    return TaskResult(False, "Email authentication required", error="Authentication failed")
            
            # Mark email as read using API
            if hasattr(client, 'mark_as_read'):
                success = client.mark_as_read(email_id)
            elif hasattr(client, '_make_request'):
                # Use outlook_integration's _make_request method to mark as read
                result = client._make_request("PATCH", f"/me/messages/{email_id}", json={"isRead": True})
                success = result is not None
            else:
                return TaskResult(False, "Email client does not support marking emails as read", error="Unsupported client interface")
            
            if success:
                return TaskResult(True, f"Email {email_id} marked as read", data={"email_id": email_id})
            else:
                return TaskResult(False, "Failed to mark email as read", error="Mark read failed")
        except ImportError as ie:
            return TaskResult(False, f"Email integration not available. Missing required packages.\n\nInstall with: pip install msal requests python-dotenv\n\nError: {str(ie)}", error="Import error")
        except Exception as e:
            logger.error(f"EmailMarkReadTask error: {e}")
            return TaskResult(False, f"Failed to mark email as read: {str(e)}", error=str(e))


# =====================================================
# DOCUMENT CREATION TASKS
# =====================================================

class WordDocumentCreateTask(BaseTask):
    """Create a Word document from text content"""
    
    def __init__(self):
        super().__init__(
            name="word_document_create",
            description="Create a Word document (.docx) with specified content. Can create new documents or append to existing ones.",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["file_path", "content"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from docx import Document
            from docx.shared import Pt
            from pathlib import Path
            
            file_path = kwargs.get("file_path")
            content = kwargs.get("content")
            title = kwargs.get("title", "")
            
            if not file_path or not content:
                return TaskResult(False, "file_path and content parameters are required", error="Missing parameters")
            
            # Ensure .docx extension
            if not file_path.endswith('.docx'):
                file_path += '.docx'
            
            file_path_obj = Path(file_path)
            file_path_obj.parent.mkdir(parents=True, exist_ok=True)
            
            # Check if file exists - if so, append to it
            if file_path_obj.exists():
                doc = Document(file_path)
                doc.add_paragraph()  # Add spacing
            else:
                doc = Document()
                if title:
                    doc.add_heading(title, 0)
            
            # Add content (preserve paragraphs if content has newlines)
            paragraphs = content.split('\n')
            for para_text in paragraphs:
                if para_text.strip():
                    para = doc.add_paragraph(para_text.strip())
                    para.style.font.size = Pt(11)
            
            doc.save(file_path)
            
            return TaskResult(
                True,
                f"Word document created/updated: {file_path}",
                data={"file_path": str(file_path_obj), "size": file_path_obj.stat().st_size}
            )
        except ImportError:
            return TaskResult(False, "python-docx package required. Install with: pip install python-docx", error="Package not available")
        except Exception as e:
            logger.error(f"WordDocumentCreateTask error: {e}")
            logger.error(traceback.format_exc())
            return TaskResult(False, f"Failed to create Word document: {str(e)}", error=str(e))


class PdfToWordTask(BaseTask):
    """Convert PDF file to Word document with extracted text"""
    
    def __init__(self):
        super().__init__(
            name="pdf_to_word",
            description="Extract text from a PDF file and save it as a Word document (.docx). Preserves text content and basic formatting.",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["pdf_path", "output_path"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from docx import Document
            from docx.shared import Pt
            from pathlib import Path
            
            pdf_path = kwargs.get("pdf_path")
            output_path = kwargs.get("output_path")
            
            if not pdf_path or not output_path:
                return TaskResult(False, "pdf_path and output_path parameters are required", error="Missing parameters")
            
            pdf_path_obj = Path(pdf_path)
            if not pdf_path_obj.exists():
                return TaskResult(False, f"PDF file not found: {pdf_path}", error="File not found")
            
            # Ensure output has .docx extension
            if not output_path.endswith('.docx'):
                output_path += '.docx'
            
            output_path_obj = Path(output_path)
            output_path_obj.parent.mkdir(parents=True, exist_ok=True)
            
            # Try pdfplumber first (better text extraction)
            try:
                import pdfplumber
                text_content = []
                with pdfplumber.open(pdf_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text:
                            text_content.append(f"--- Page {page_num} ---\n{text}\n")
                
                extracted_text = "\n".join(text_content)
            except ImportError:
                # Fallback to pypdf
                try:
                    import pypdf
                    text_content = []
                    with open(pdf_path, 'rb') as file:
                        pdf_reader = pypdf.PdfReader(file)
                        for page_num, page in enumerate(pdf_reader.pages, 1):
                            text = page.extract_text()
                            if text:
                                text_content.append(f"--- Page {page_num} ---\n{text}\n")
                    extracted_text = "\n".join(text_content)
                except ImportError:
                    return TaskResult(False, "PDF library required. Install with: pip install pdfplumber or pip install pypdf", error="Package not available")
            
            if not extracted_text.strip():
                return TaskResult(False, "No text could be extracted from the PDF. The PDF may be image-based or encrypted.", error="No text extracted")
            
            # Create Word document
            doc = Document()
            doc.add_heading(f"Extracted from: {pdf_path_obj.name}", 0)
            
            # Add extracted text
            paragraphs = extracted_text.split('\n')
            for para_text in paragraphs:
                if para_text.strip():
                    para = doc.add_paragraph(para_text.strip())
                    para.style.font.size = Pt(11)
            
            doc.save(output_path)
            
            return TaskResult(
                True,
                f"PDF converted to Word document: {output_path}",
                data={
                    "pdf_path": str(pdf_path_obj),
                    "output_path": str(output_path_obj),
                    "output_size": output_path_obj.stat().st_size,
                    "pages_extracted": len(text_content)
                }
            )
        except Exception as e:
            logger.error(f"PdfToWordTask error: {e}")
            logger.error(traceback.format_exc())
            return TaskResult(False, f"Failed to convert PDF to Word: {str(e)}", error=str(e))


class PowerPointFormatTextTask(BaseTask):
    """Format text in PowerPoint presentations"""
    
    def __init__(self):
        super().__init__(
            name="powerpoint_format_text",
            description="Format text in PowerPoint presentations (.pptx). Can change font properties (size, bold, italic, color, font name), alignment, and replace text content in specific slides or all slides.",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["file_path"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            from pathlib import Path
            
            file_path = kwargs.get("file_path")
            if not file_path:
                return TaskResult(False, "file_path parameter is required", error="Missing parameter")
            
            # Ensure .pptx extension
            if not file_path.endswith('.pptx'):
                file_path += '.pptx'
            
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                return TaskResult(False, f"PowerPoint file not found: {file_path}", error="File not found")
            
            # Check if file can be written (basic permission check)
            try:
                # Try to open file for writing to check permissions
                test_file = open(file_path, 'r+b')
                test_file.close()
            except PermissionError:
                return TaskResult(False, f"Permission denied: {file_path}. The file may be open in PowerPoint or another program. Please close PowerPoint and try again.", error="Permission denied - file may be open")
            except Exception:
                # If it's not a permission error, continue (might be a different issue)
                pass
            
            # Open presentation
            prs = Presentation(file_path)
            
            # Get formatting parameters
            slide_number = kwargs.get("slide_number")  # 1-based, None means all slides
            search_text = kwargs.get("search_text")  # Text to find (optional)
            replace_text = kwargs.get("replace_text")  # Text to replace with (optional)
            
            # Case fixing options
            fix_case = kwargs.get("fix_case", False)  # Boolean - enable automatic case fixing
            title_case_for_titles = kwargs.get("title_case_for_titles", True)  # Apply title case to titles
            sentence_case_for_body = kwargs.get("sentence_case_for_body", True)  # Apply sentence case to body
            
            # Font formatting options
            font_size = kwargs.get("font_size")  # Integer, in points
            font_bold = kwargs.get("font_bold")  # Boolean
            font_italic = kwargs.get("font_italic")  # Boolean
            font_color = kwargs.get("font_color")  # Hex color string like "#FF0000" or RGB tuple (255, 0, 0)
            font_name = kwargs.get("font_name")  # Font family name
            
            # Alignment
            alignment = kwargs.get("alignment")  # "left", "center", "right", "justify"
            
            # Helper function for title case (capitalize first letter of each word, handle small words)
            def to_title_case(text):
                """Convert text to title case, handling small words appropriately"""
                if not text:
                    return text
                # List of small words that should be lowercase unless first word
                small_words = {'a', 'an', 'and', 'as', 'at', 'but', 'by', 'for', 'from', 
                              'in', 'into', 'nor', 'of', 'on', 'or', 'the', 'to', 'with'}
                words = text.split()
                if not words:
                    return text
                # Always capitalize first word
                result = [words[0].capitalize()]
                for word in words[1:]:
                    if word.lower() in small_words:
                        result.append(word.lower())
                    else:
                        result.append(word.capitalize())
                return ' '.join(result)
            
            # Helper function for sentence case (capitalize first letter of sentence)
            def to_sentence_case(text):
                """Convert text to sentence case"""
                if not text:
                    return text
                # Split by sentences (period, exclamation, question mark)
                import re
                sentences = re.split(r'([.!?]\s+)', text)
                result = []
                for i, part in enumerate(sentences):
                    if i % 2 == 0:  # The sentence text
                        if part.strip():
                            result.append(part[0].upper() + part[1:].lower() if len(part) > 1 else part.upper())
                        else:
                            result.append(part)
                    else:  # The punctuation and space
                        result.append(part)
                return ''.join(result)
            
            # Determine which slides to process
            if slide_number is not None:
                slide_index = int(slide_number) - 1  # Convert to 0-based
                if slide_index < 0 or slide_index >= len(prs.slides):
                    return TaskResult(False, f"Slide number {slide_number} is out of range. Presentation has {len(prs.slides)} slides.", error="Invalid slide number")
                slides_to_process = [prs.slides[slide_index]]
            else:
                slides_to_process = prs.slides
            
            changes_made = 0
            slides_modified = 0
            
            # Process each slide
            for slide in slides_to_process:
                slide_modified = False
                
                # Process all shapes on the slide
                for shape in slide.shapes:
                    if not hasattr(shape, "text_frame"):
                        continue
                    
                    # Detect if this is a title placeholder (usually first shape or has placeholder type 1)
                    is_title = False
                    try:
                        if hasattr(shape, "placeholder_format"):
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                is_title = True
                        # Also check if it's the first text shape (common title position)
                        if shape == slide.shapes[0] and hasattr(shape, "text_frame"):
                            is_title = True
                    except:
                        pass
                    
                    # Process all paragraphs in the text frame
                    for paragraph in shape.text_frame.paragraphs:
                        # Check if we need to replace text
                        if search_text and replace_text and search_text in paragraph.text:
                            paragraph.text = paragraph.text.replace(search_text, replace_text)
                            slide_modified = True
                            changes_made += 1
                        
                        # Apply case fixing if enabled
                        if fix_case:
                            original_text = paragraph.text
                            if original_text and original_text.strip():
                                if is_title and title_case_for_titles:
                                    new_text = to_title_case(original_text)
                                    if new_text != original_text:
                                        paragraph.text = new_text
                                        slide_modified = True
                                        changes_made += 1
                                elif not is_title and sentence_case_for_body:
                                    new_text = to_sentence_case(original_text)
                                    if new_text != original_text:
                                        paragraph.text = new_text
                                        slide_modified = True
                                        changes_made += 1
                        
                        # Process all runs in the paragraph
                        for run in paragraph.runs:
                            text_to_format = run.text
                            
                            # If search_text is specified, only format matching runs
                            if search_text and search_text not in text_to_format:
                                continue
                            
                            # Apply font formatting
                            if font_size is not None:
                                try:
                                    run.font.size = Pt(int(font_size))
                                    slide_modified = True
                                    changes_made += 1
                                except (ValueError, TypeError):
                                    pass
                            
                            if font_bold is not None:
                                run.font.bold = bool(font_bold)
                                slide_modified = True
                                changes_made += 1
                            
                            if font_italic is not None:
                                run.font.italic = bool(font_italic)
                                slide_modified = True
                                changes_made += 1
                            
                            if font_name is not None:
                                run.font.name = str(font_name)
                                slide_modified = True
                                changes_made += 1
                            
                            if font_color is not None:
                                try:
                                    # Handle hex color string
                                    if isinstance(font_color, str) and font_color.startswith('#'):
                                        hex_color = font_color.lstrip('#')
                                        r = int(hex_color[0:2], 16)
                                        g = int(hex_color[2:4], 16)
                                        b = int(hex_color[4:6], 16)
                                        run.font.color.rgb = RGBColor(r, g, b)
                                    # Handle RGB tuple
                                    elif isinstance(font_color, (list, tuple)) and len(font_color) == 3:
                                        r, g, b = font_color
                                        run.font.color.rgb = RGBColor(int(r), int(g), int(b))
                                    else:
                                        # Try to parse as hex without #
                                        hex_color = str(font_color).lstrip('#')
                                        if len(hex_color) == 6:
                                            r = int(hex_color[0:2], 16)
                                            g = int(hex_color[2:4], 16)
                                            b = int(hex_color[4:6], 16)
                                            run.font.color.rgb = RGBColor(r, g, b)
                                    slide_modified = True
                                    changes_made += 1
                                except Exception as color_error:
                                    logger.warning(f"Error setting font color: {color_error}")
                            
                            # Apply alignment
                            if alignment is not None:
                                try:
                                    align_map = {
                                        "left": PP_ALIGN.LEFT,
                                        "center": PP_ALIGN.CENTER,
                                        "right": PP_ALIGN.RIGHT,
                                        "justify": PP_ALIGN.JUSTIFY
                                    }
                                    if alignment.lower() in align_map:
                                        paragraph.alignment = align_map[alignment.lower()]
                                        slide_modified = True
                                        changes_made += 1
                                except Exception as align_error:
                                    logger.warning(f"Error setting alignment: {align_error}")
                    
                    # Also check for table shapes
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        # Apply same formatting to table cells
                                        if font_size is not None:
                                            try:
                                                run.font.size = Pt(int(font_size))
                                                slide_modified = True
                                                changes_made += 1
                                            except (ValueError, TypeError):
                                                pass
                                        if font_bold is not None:
                                            run.font.bold = bool(font_bold)
                                            slide_modified = True
                                            changes_made += 1
                                        if font_italic is not None:
                                            run.font.italic = bool(font_italic)
                                            slide_modified = True
                                            changes_made += 1
                                        if font_name is not None:
                                            run.font.name = str(font_name)
                                            slide_modified = True
                                            changes_made += 1
                
                if slide_modified:
                    slides_modified += 1
            
            # Save the presentation
            try:
                prs.save(file_path)
            except PermissionError as save_error:
                return TaskResult(False, f"Permission denied when saving: {file_path}. The file may be open in PowerPoint or another program. Please close PowerPoint and try again.", error=f"Permission denied: {str(save_error)}")
            except IOError as io_error:
                return TaskResult(False, f"Error saving file: {file_path}. {str(io_error)}. Make sure the file is not open in PowerPoint.", error=f"IO error: {str(io_error)}")
            
            # Build result message
            result_parts = [f"PowerPoint presentation updated: {file_path_obj.name}"]
            if slide_number:
                result_parts.append(f"Modified slide {slide_number}")
            else:
                result_parts.append(f"Modified {slides_modified} slide(s)")
            
            if changes_made > 0:
                result_parts.append(f"Applied {changes_made} formatting change(s)")
            else:
                result_parts.append("No changes were needed - text may already be correctly formatted, or text may be in images/shapes that can't be automatically edited")
            
            if fix_case:
                if title_case_for_titles:
                    result_parts.append("Applied title case to titles")
                if sentence_case_for_body:
                    result_parts.append("Applied sentence case to body text")
            
            if search_text and replace_text:
                result_parts.append(f"Replaced '{search_text}' with '{replace_text}'")
            
            return TaskResult(
                True,
                ". ".join(result_parts) + ".",
                data={
                    "file_path": str(file_path_obj),
                    "slides_modified": slides_modified,
                    "changes_made": changes_made,
                    "slide_number": slide_number
                }
            )
        except ImportError:
            return TaskResult(False, "python-pptx package required. Install with: pip install python-pptx", error="Package not available")
        except Exception as e:
            logger.error(f"PowerPointFormatTextTask error: {e}")
            logger.error(traceback.format_exc())
            return TaskResult(False, f"Failed to format PowerPoint: {str(e)}", error=str(e))


class FileViewTask(BaseTask):
    """View/read content from various file types (text, PDF, images, etc.)"""
    
    def __init__(self):
        super().__init__(
            name="file_view",
            description="View and extract content from various file types including text files, PDFs, images (with OCR), Word documents, Excel files, etc. Returns file content or summary.",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["file_path"]
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from pathlib import Path
            
            file_path = kwargs.get("file_path")
            if not file_path:
                return TaskResult(False, "file_path parameter is required", error="Missing parameter")
            
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                return TaskResult(False, f"File not found: {file_path}", error="File not found")
            
            file_ext = file_path_obj.suffix.lower()
            content_summary = ""
            file_info = {
                "path": str(file_path_obj),
                "size": file_path_obj.stat().st_size,
                "extension": file_ext
            }
            
            # Handle different file types
            if file_ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']:
                # Text files
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    content_summary = f"Text file content ({len(content)} characters):\n\n{content[:5000]}"  # First 5000 chars
                    if len(content) > 5000:
                        content_summary += f"\n\n... (truncated, total {len(content)} characters)"
                    file_info["content"] = content
                    file_info["char_count"] = len(content)
                except Exception as e:
                    content_summary = f"Error reading text file: {str(e)}"
            
            elif file_ext == '.pdf':
                # PDF files
                try:
                    import pdfplumber
                    text_parts = []
                    with pdfplumber.open(file_path) as pdf:
                        for page_num, page in enumerate(pdf.pages[:5], 1):  # First 5 pages
                            text = page.extract_text()
                            if text:
                                text_parts.append(f"Page {page_num}:\n{text[:1000]}")  # First 1000 chars per page
                    content_summary = "PDF content (first 5 pages):\n\n" + "\n\n".join(text_parts)
                    file_info["pages"] = len(pdf.pages) if 'pdf' in locals() else 0
                except ImportError:
                    try:
                        import pypdf
                        text_parts = []
                        with open(file_path, 'rb') as file:
                            pdf_reader = pypdf.PdfReader(file)
                            for page_num, page in enumerate(pdf_reader.pages[:5], 1):
                                text = page.extract_text()
                                if text:
                                    text_parts.append(f"Page {page_num}:\n{text[:1000]}")
                        content_summary = "PDF content (first 5 pages):\n\n" + "\n\n".join(text_parts)
                        file_info["pages"] = len(pdf_reader.pages)
                    except ImportError:
                        content_summary = "PDF file detected but no PDF library available. Install pdfplumber or pypdf."
                except Exception as e:
                    content_summary = f"Error reading PDF: {str(e)}"
            
            elif file_ext in ['.docx', '.doc']:
                # Word documents
                try:
                    from docx import Document
                    doc = Document(file_path)
                    paragraphs = [para.text for para in doc.paragraphs[:50]]  # First 50 paragraphs
                    content_summary = "Word document content:\n\n" + "\n".join(paragraphs)
                    if len(doc.paragraphs) > 50:
                        content_summary += f"\n\n... (truncated, total {len(doc.paragraphs)} paragraphs)"
                    file_info["paragraphs"] = len(doc.paragraphs)
                except ImportError:
                    content_summary = "Word document detected but python-docx not available. Install with: pip install python-docx"
                except Exception as e:
                    content_summary = f"Error reading Word document: {str(e)}"
            
            elif file_ext in ['.xlsx', '.xls']:
                # Excel files - viewing not available (pandas removed)
                content_summary = f"Excel file detected ({file_path_obj.name}). File size: {file_path_obj.stat().st_size} bytes.\n\nNote: Excel file viewing is not available. You can open this file in Excel or ask me to convert it to CSV for viewing."
            
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                # Image files - try OCR
                try:
                    import pytesseract
                    from PIL import Image
                    img = Image.open(file_path)
                    text = pytesseract.image_to_string(img)
                    if text.strip():
                        content_summary = f"Image file - OCR extracted text:\n\n{text[:2000]}"
                        file_info["ocr_text"] = text
                    else:
                        content_summary = "Image file - no text detected via OCR"
                except ImportError:
                    content_summary = f"Image file detected ({file_path_obj.name}). OCR not available. Install pytesseract and tesseract-ocr for text extraction."
                except Exception as e:
                    content_summary = f"Image file detected but OCR error: {str(e)}"
            
            else:
                content_summary = f"File type {file_ext} detected. File size: {file_path_obj.stat().st_size} bytes. Content viewing not supported for this file type."
            
            return TaskResult(
                True,
                content_summary,
                data=file_info
            )
        except Exception as e:
            logger.error(f"FileViewTask error: {e}")
            logger.error(traceback.format_exc())
            return TaskResult(False, f"Failed to view file: {str(e)}", error=str(e))
            
            client = get_email_client()
            if not client:
                return TaskResult(False, "Email client not available", error="Email client not initialized")
            
            if not client.is_authenticated():
                if not client.authenticate():
                    return TaskResult(False, "Email authentication required", error="Authentication failed")
            
            success = client.mark_as_read(email_id)
            
            if success:
                return TaskResult(True, f"Email {email_id} marked as read", data={"email_id": email_id})
            else:
                return TaskResult(False, "Failed to mark email as read", error="Mark read failed")
        except ImportError:
            return TaskResult(False, "Email integration not available", error="Import error")
        except Exception as e:
            logger.error(f"EmailMarkReadTask error: {e}")
            return TaskResult(False, f"Failed to mark email as read: {str(e)}", error=str(e))


# =====================================================
# TASK REGISTRY
# =====================================================

class TaskRegistry:
    """Registry for all available tasks"""
    
    def __init__(self):
        self.tasks: Dict[str, BaseTask] = {}
        self.task_history: List[Dict] = []
        self.config_path = Path("lea_tasks_config.json")
        self.load_config()
    
    def register_task(self, task: BaseTask):
        """Register a task"""
        self.tasks[task.name] = task
        logger.info(f"Registered task: {task.name}")
    
    def get_task(self, task_name: str) -> Optional[BaseTask]:
        """Get a task by name"""
        return self.tasks.get(task_name)
    
    def list_tasks(self) -> List[Dict]:
        """List all available tasks"""
        return [
            {
                "name": task.name,
                "description": task.description,
                "requires_confirmation": task.requires_confirmation,
                "allowed": task.allowed,
                "required_params": task.get_required_params()
            }
            for task in self.tasks.values()
        ]
    
    def execute_task(self, task_name: str, params: Dict, confirmed: bool = False) -> TaskResult:
        """Execute a task"""
        task = self.get_task(task_name)
        
        if not task:
            return TaskResult(False, f"Task not found: {task_name}", error="Task not found")
        
        if not task.allowed:
            return TaskResult(False, f"Task is disabled: {task_name}", error="Task disabled")
        
        if task.requires_confirmation and not confirmed:
            return TaskResult(False, f"Task requires confirmation: {task_name}", error="Confirmation required")
        
        # Execute
        result = task.execute(**params)
        
        # Log to history
        self.task_history.append({
            "timestamp": datetime.now().isoformat(),
            "task_name": task_name,
            "params": params,
            "result": result.to_dict(),
            "confirmed": confirmed
        })
        
        # Save history (keep last 100)
        if len(self.task_history) > 100:
            self.task_history = self.task_history[-100:]
        
        self.save_config()
        
        return result
    
    def load_config(self):
        """Load task configuration"""
        if self.config_path.exists():
            try:
                with open(self.config_path, 'r') as f:
                    config = json.load(f)
                    # Apply config to tasks
                    for task_name, task_config in config.get("tasks", {}).items():
                        if task_name in self.tasks:
                            self.tasks[task_name].allowed = task_config.get("allowed", True)
            except Exception as e:
                logger.error(f"Error loading task config: {e}")
    
    def save_config(self):
        """Save task configuration"""
        try:
            config = {
                "tasks": {
                    name: {"allowed": task.allowed}
                    for name, task in self.tasks.items()
                },
                "history": self.task_history[-50:]  # Keep last 50
            }
            with open(self.config_path, 'w') as f:
                json.dump(config, f, indent=2)
        except Exception as e:
            logger.error(f"Error saving task config: {e}")
    
    def enable_task(self, task_name: str):
        """Enable a task"""
        if task_name in self.tasks:
            self.tasks[task_name].allowed = True
            self.save_config()
    
    def disable_task(self, task_name: str):
        """Disable a task"""
        if task_name in self.tasks:
            self.tasks[task_name].allowed = False
            self.save_config()


# =====================================================
# WORKFLOW AUTOMATION TASKS (Executive Assistant Only)
# =====================================================

class WorkflowRecordTask(BaseTask):
    """Record a new workflow by watching user actions"""
    
    def __init__(self):
        super().__init__(
            name="workflow_record",
            description="Start recording a new workflow. Lea will watch and record your actions until you stop recording.",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["workflow_name"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        workflow_name = kwargs.get("workflow_name")
        if not workflow_name:
            return False, "Workflow name is required"
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from workflow_system import get_workflow_manager
            workflow_name = kwargs.get("workflow_name")
            description = kwargs.get("description", "")
            
            manager = get_workflow_manager()
            success, message = manager.recorder.start_recording(workflow_name)
            
            if success:
                return TaskResult(
                    True,
                    f"Recording started for workflow '{workflow_name}'. Perform the actions you want to record, then use 'workflow_stop' to finish.",
                    data={"workflow_name": workflow_name, "status": "recording"}
                )
            else:
                return TaskResult(False, message, error=message)
        except ImportError:
            return TaskResult(False, "Workflow system not available", error="workflow_system module not found")
        except Exception as e:
            return TaskResult(False, f"Error starting workflow recording: {str(e)}", error=str(e))


class WorkflowStopTask(BaseTask):
    """Stop recording current workflow and save it"""
    
    def __init__(self):
        super().__init__(
            name="workflow_stop",
            description="Stop recording the current workflow and save it for future use.",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["workflow_name", "description"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        workflow_name = kwargs.get("workflow_name")
        description = kwargs.get("description")
        if not workflow_name:
            return False, "Workflow name is required"
        if not description:
            return False, "Workflow description is required"
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from workflow_system import get_workflow_manager, Workflow
            from datetime import datetime
            
            workflow_name = kwargs.get("workflow_name")
            description = kwargs.get("description", "")
            parameters = kwargs.get("parameters", {})  # Dict of parameter names and descriptions
            
            manager = get_workflow_manager()
            success, message, actions = manager.recorder.stop_recording()
            
            if not success:
                return TaskResult(False, message, error=message)
            
            if not actions:
                return TaskResult(False, "No actions were recorded", error="No actions recorded")
            
            # Create workflow object
            workflow = Workflow(
                name=workflow_name,
                description=description,
                created=datetime.now().isoformat(),
                modified=datetime.now().isoformat(),
                actions=actions,
                parameters=parameters if isinstance(parameters, dict) else {},
                category=kwargs.get("category", "general")
            )
            
            # Save workflow
            save_success, save_message = manager.save_workflow(workflow)
            
            if save_success:
                return TaskResult(
                    True,
                    f"Workflow '{workflow_name}' saved successfully with {len(actions)} actions.",
                    data={"workflow_name": workflow_name, "action_count": len(actions)}
                )
            else:
                return TaskResult(False, save_message, error=save_message)
        
        except ImportError:
            return TaskResult(False, "Workflow system not available", error="workflow_system module not found")
        except Exception as e:
            return TaskResult(False, f"Error stopping workflow recording: {str(e)}", error=str(e))


class WorkflowPlayTask(BaseTask):
    """Play back a saved workflow"""
    
    def __init__(self):
        super().__init__(
            name="workflow_play",
            description="Execute a saved workflow. Can include parameters to customize the workflow execution.",
            requires_confirmation=False
        )
    
    def get_required_params(self) -> List[str]:
        return ["workflow_name"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        workflow_name = kwargs.get("workflow_name")
        if not workflow_name:
            return False, "Workflow name is required"
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from workflow_system import get_workflow_manager
            
            workflow_name = kwargs.get("workflow_name")
            parameters = kwargs.get("parameters", {})  # Parameters to substitute in workflow
            
            manager = get_workflow_manager()
            success, message, workflow = manager.load_workflow(workflow_name)
            
            if not success or not workflow:
                return TaskResult(False, message or f"Workflow '{workflow_name}' not found", error=message)
            
            # Play the workflow
            play_success, play_message = manager.player.play_workflow(workflow, parameters)
            
            if play_success:
                return TaskResult(
                    True,
                    f"Workflow '{workflow_name}' executed successfully. {play_message}",
                    data={"workflow_name": workflow_name, "parameters": parameters}
                )
            else:
                return TaskResult(False, play_message, error=play_message)
        
        except ImportError:
            return TaskResult(False, "Workflow system not available", error="workflow_system module not found")
        except Exception as e:
            return TaskResult(False, f"Error playing workflow: {str(e)}", error=str(e))


class WorkflowListTask(BaseTask):
    """List all available workflows"""
    
    def __init__(self):
        super().__init__(
            name="workflow_list",
            description="List all saved workflows with their descriptions and parameters.",
            requires_confirmation=False
        )
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from workflow_system import get_workflow_manager
            
            manager = get_workflow_manager()
            workflows = manager.list_workflows()
            
            if not workflows:
                return TaskResult(
                    True,
                    "No workflows found. Use 'workflow_record' to create a new workflow.",
                    data={"workflows": []}
                )
            
            # Format workflow list for display
            workflow_list = []
            for wf in workflows:
                workflow_list.append({
                    "name": wf["name"],
                    "description": wf.get("description", ""),
                    "category": wf.get("category", "general"),
                    "action_count": wf.get("action_count", 0),
                    "parameters": wf.get("parameters", {}),
                    "created": wf.get("created", ""),
                    "modified": wf.get("modified", "")
                })
            
            return TaskResult(
                True,
                f"Found {len(workflows)} workflow(s).",
                data={"workflows": workflow_list, "count": len(workflows)}
            )
        
        except ImportError:
            return TaskResult(False, "Workflow system not available", error="workflow_system module not found")
        except Exception as e:
            return TaskResult(False, f"Error listing workflows: {str(e)}", error=str(e))


class WorkflowDeleteTask(BaseTask):
    """Delete a saved workflow"""
    
    def __init__(self):
        super().__init__(
            name="workflow_delete",
            description="Delete a saved workflow. This action cannot be undone.",
            requires_confirmation=True  # Require confirmation for deletion
        )
    
    def get_required_params(self) -> List[str]:
        return ["workflow_name"]
    
    def validate_params(self, **kwargs) -> Tuple[bool, str]:
        workflow_name = kwargs.get("workflow_name")
        if not workflow_name:
            return False, "Workflow name is required"
        return True, ""
    
    def execute(self, **kwargs) -> TaskResult:
        try:
            from workflow_system import get_workflow_manager
            
            workflow_name = kwargs.get("workflow_name")
            manager = get_workflow_manager()
            success, message = manager.delete_workflow(workflow_name)
            
            if success:
                return TaskResult(True, message, data={"workflow_name": workflow_name})
            else:
                return TaskResult(False, message, error=message)
        
        except ImportError:
            return TaskResult(False, "Workflow system not available", error="workflow_system module not found")
        except Exception as e:
            return TaskResult(False, f"Error deleting workflow: {str(e)}", error=str(e))


# =====================================================
# INITIALIZE REGISTRY WITH BUILT-IN TASKS
# =====================================================

def create_task_registry() -> TaskRegistry:
    """Create and populate task registry"""
    registry = TaskRegistry()
    
    # Register all built-in tasks
    registry.register_task(FileCopyTask())
    registry.register_task(FileMoveTask())
    registry.register_task(FileDeleteTask())
    registry.register_task(FileReadTask())
    registry.register_task(FileWriteTask())
    registry.register_task(FileViewTask())  # New: View various file types
    registry.register_task(DirectoryCreateTask())
    registry.register_task(DirectoryListTask())
    registry.register_task(TextReplaceTask())
    registry.register_task(SystemCommandTask())
    
    # Register document creation tasks
    registry.register_task(WordDocumentCreateTask())
    registry.register_task(PdfToWordTask())
    registry.register_task(PowerPointFormatTextTask())
    
    # Register screen automation tasks
    registry.register_task(ScreenshotTask())
    registry.register_task(ClickTask())
    registry.register_task(TypeTask())
    registry.register_task(KeyPressTask())
    registry.register_task(HotkeyTask())
    registry.register_task(FindImageTask())
    registry.register_task(ScrollTask())
    registry.register_task(MoveMouseTask())
    registry.register_task(GetScreenSizeTask())
    
    # Register workflow automation tasks (Executive Assistant mode only)
    try:
        # Try to import workflow system - if it fails, just skip workflow tasks
        try:
            from workflow_system import get_workflow_manager
            # Test if workflow system is available (this might fail if dependencies missing)
            try:
                manager = get_workflow_manager()
                # Only register if we got this far
                registry.register_task(WorkflowRecordTask())
                registry.register_task(WorkflowStopTask())
                registry.register_task(WorkflowPlayTask())
                registry.register_task(WorkflowListTask())
                registry.register_task(WorkflowDeleteTask())
                logger.info("Workflow automation tasks registered successfully")
            except Exception as manager_error:
                logger.warning(f"Workflow manager initialization failed: {manager_error} - workflow tasks not registered")
        except ImportError as import_error:
            logger.warning(f"Workflow system module not found: {import_error} - workflow tasks not registered")
    except Exception as e:
        # Catch-all to prevent workflow registration from breaking the entire task system
        logger.warning(f"Error registering workflow tasks (non-critical): {e}")
        logger.warning("Continuing without workflow tasks - other tasks will still work")
    
    # Register Outlook tasks (new naming convention)
    try:
        # Try to import outlook integration to check if it's available
        try:
            from outlook_integration import OutlookClient
            client_available = True
        except ImportError:
            try:
                from outlook_email import get_email_client
                client_available = get_email_client() is not None
            except ImportError:
                client_available = False
        
        if client_available:
            # Register new Outlook tasks with correct names
            registry.register_task(OutlookEmailCheckTask())
            registry.register_task(OutlookEmailDraftTask())
            registry.register_task(OutlookCalendarCheckTask())
            registry.register_task(OutlookInboxOrganizeTask())
            registry.register_task(OutlookUserProfileTask())
            logger.info("Outlook tasks registered successfully")
            
            # Also register legacy email tasks for backward compatibility
            registry.register_task(EmailCheckTask())
            registry.register_task(EmailReadTask())
            registry.register_task(EmailSendTask())
            registry.register_task(EmailMarkReadTask())
            logger.info("Legacy email tasks registered for backward compatibility")
        else:
            # Register tasks anyway - they'll handle errors gracefully
            registry.register_task(OutlookEmailCheckTask())
            registry.register_task(OutlookEmailDraftTask())
            registry.register_task(OutlookCalendarCheckTask())
            registry.register_task(OutlookInboxOrganizeTask())
            registry.register_task(OutlookUserProfileTask())
            logger.info("Outlook tasks registered (will show error if Outlook not configured)")
    except Exception as e:
        logger.warning(f"Error registering Outlook tasks: {e}")
        # Register tasks anyway - they'll handle errors gracefully
        try:
            registry.register_task(OutlookEmailCheckTask())
            registry.register_task(OutlookEmailDraftTask())
            registry.register_task(OutlookCalendarCheckTask())
            registry.register_task(OutlookInboxOrganizeTask())
            registry.register_task(OutlookUserProfileTask())
            logger.info("Outlook tasks registered (will show error if Outlook not configured)")
        except Exception as e2:
            logger.error(f"Failed to register Outlook tasks: {e2}")
    
    # Try to register custom tasks (if file exists)
    try:
        # Try importing custom tasks module
        import importlib.util
        custom_tasks_path = Path(__file__).parent / "custom_tasks_example.py"
        
        if custom_tasks_path.exists():
            spec = importlib.util.spec_from_file_location("custom_tasks", custom_tasks_path)
            if spec and spec.loader:
                custom_tasks = importlib.util.module_from_spec(spec)
                spec.loader.exec_module(custom_tasks)
                
                # Call register function if it exists
                if hasattr(custom_tasks, 'register_custom_tasks'):
                    custom_tasks.register_custom_tasks()
                    logging.info("Custom tasks loaded successfully")
    except Exception as e:
        # Custom tasks are optional, just log and continue
        logging.info(f"Custom tasks not loaded (this is OK): {e}")
    
    return registry


# Global registry instance
_task_registry = None

def get_task_registry() -> TaskRegistry:
    """Get or create the global task registry"""
    global _task_registry
    if _task_registry is None:
        _task_registry = create_task_registry()
    return _task_registry

